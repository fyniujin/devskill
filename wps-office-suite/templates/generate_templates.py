"""
纯 Python 代码生成模板（无任何二进制文件）
运行此脚本即可生成 3 个常用模板文件

用法：
  python templates/generate_templates.py          # 生成到当前工作目录
  python templates/generate_templates.py --dir ~   # 生成到指定目录
"""
import sys
from pathlib import Path

def generate_report_template(output_dir: Path):
    """生成工作报告模板（周报/月报/汇报）"""
    try:
        from docx import Document
        from docx.shared import Pt, Cm, RGBColor
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        
        doc = Document()
        
        # 标题
        title = doc.add_heading("工作报告", level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        # 基本信息
        doc.add_paragraph("报告周期：____年____月____日 至 ____年____月____日")
        p = doc.add_paragraph()
        p.add_run("报告人：").bold = True
        p.add_run("____________")
        p = doc.add_paragraph()
        p.add_run("部门：").bold = True
        p.add_run("____________")
        
        doc.add_paragraph("")
        
        # 本周工作完成情况
        doc.add_heading("一、本周工作完成情况", level=1)
        doc.add_paragraph("（请按重要性排序，每项简要说明）")
        for i in range(1, 6):
            p = doc.add_paragraph(f"{i}. ", style='List Number')
            p.add_run("________________________")
        
        doc.add_paragraph("")
        
        # 遇到的问题
        doc.add_heading("二、遇到的问题与风险", level=1)
        doc.add_paragraph("（描述当前阻塞项、延期风险及需要协助的事项）")
        for i in range(1, 4):
            p = doc.add_paragraph(f"{i}. ", style='List Number')
            p.add_run("________________________")
        
        doc.add_paragraph("")
        
        # 下周计划
        doc.add_heading("三、下周工作计划", level=1)
        doc.add_paragraph("（列出下周重点任务，注明预计完成时间）")
        for i in range(1, 5):
            p = doc.add_paragraph(f"{i}. ", style='List Number')
            p.add_run("________________________")
        
        doc.add_paragraph("")
        
        # 备注
        doc.add_heading("四、备注", level=1)
        doc.add_paragraph("（其他需要说明的事项）")
        doc.add_paragraph("___________________________________________________")
        
        # 保存
        output_dir.mkdir(parents=True, exist_ok=True)
        out = output_dir / "report_template.docx"
        doc.save(out)
        return str(out)
    except ImportError:
        return None


def generate_budget_template(output_dir: Path):
    """生成月度预算模板（收入/支出/汇总）"""
    try:
        import openpyxl
        from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
        
        wb = openpyxl.Workbook()
        wb.remove(wb.active)
        
        # ===== 收入 Sheet =====
        ws_income = wb.create_sheet("收入")
        headers = ["日期", "类别", "金额", "备注"]
        income_data = [
            ["2026-01-05", "工资收入", 12000, "月工资"],
            ["2026-01-10", "兼职收入", 3000, "副业"],
            ["2026-01-15", "理财收益", 500, "基金分红"],
            ["2026-01-20", "其他收入", 200, "红包"],
        ]
        
        header_font = Font(bold=True, color="FFFFFF", size=11)
        header_fill = PatternFill("solid", fgColor="27AE60")
        header_align = Alignment(horizontal="center", vertical="center")
        
        for c, h in enumerate(headers, 1):
            cell = ws_income.cell(row=1, column=c, value=h)
            cell.font = header_font
            cell.fill = header_fill
            cell.alignment = header_align
        
        for r, row in enumerate(income_data, 2):
            for c, val in enumerate(row, 1):
                ws_income.cell(row=r, column=c, value=val)
        
        ws_income.column_dimensions['A'].width = 12
        ws_income.column_dimensions['B'].width = 14
        ws_income.column_dimensions['C'].width = 12
        ws_income.column_dimensions['D'].width = 20
        
        # ===== 支出 Sheet =====
        ws_expense = wb.create_sheet("支出")
        expense_headers = ["日期", "类别", "金额", "备注"]
        expense_data = [
            ["2026-01-03", "餐饮", 850, "日常用餐"],
            ["2026-01-05", "交通", 320, "地铁+打车"],
            ["2026-01-08", "房租", 3500, "月租金"],
            ["2026-01-12", "购物", 1200, "日用品"],
            ["2026-01-18", "娱乐", 400, "聚餐+电影"],
            ["2026-01-25", "医疗", 200, "感冒药"],
        ]
        
        expense_fill = PatternFill("solid", fgColor="E74C3C")
        
        for c, h in enumerate(expense_headers, 1):
            cell = ws_expense.cell(row=1, column=c, value=h)
            cell.font = header_font
            cell.fill = expense_fill
            cell.alignment = header_align
        
        for r, row in enumerate(expense_data, 2):
            for c, val in enumerate(row, 1):
                ws_expense.cell(row=r, column=c, value=val)
        
        ws_expense.column_dimensions['A'].width = 12
        ws_expense.column_dimensions['B'].width = 10
        ws_expense.column_dimensions['C'].width = 10
        ws_expense.column_dimensions['D'].width = 20
        
        # ===== 汇总 Sheet =====
        ws_summary = wb.create_sheet("汇总")
        summary_headers = ["项目", "金额(元)"]
        summary_data = [
            ["总收入", "=SUM(收入!C:C)"],
            ["工资收入", '=SUMIF(收入!B:B,"工资收入",收入!C:C)'],
            ["兼职收入", '=SUMIF(收入!B:B,"兼职收入",收入!C:C)'],
            ["理财收益", '=SUMIF(收入!B:B,"理财收益",收入!C:C)'],
            ["总支出", "=SUM(支出!C:C)"],
            ["餐饮支出", '=SUMIF(支出!B:B,"餐饮",支出!C:C)'],
            ["交通支出", '=SUMIF(支出!B:B,"交通",支出!C:C)'],
            ["房租支出", '=SUMIF(支出!B:B,"房租",支出!C:C)'],
            ["结余", "=B1-B5"],
        ]
        
        summary_fill = PatternFill("solid", fgColor="3498DB")
        
        for c, h in enumerate(summary_headers, 1):
            cell = ws_summary.cell(row=1, column=c, value=h)
            cell.font = header_font
            cell.fill = summary_fill
            cell.alignment = header_align
        
        for r, row in enumerate(summary_data, 2):
            for c, val in enumerate(row, 1):
                ws_summary.cell(row=r, column=c, value=val)
        
        # 高亮结余行
        last_row = len(summary_data) + 1
        for c in range(1, 3):
            cell = ws_summary.cell(row=last_row, column=c)
            cell.font = Font(bold=True, size=12)
            cell.fill = PatternFill("solid", fgColor="F39C12")
        
        ws_summary.column_dimensions['A'].width = 14
        ws_summary.column_dimensions['B'].width = 16
        
        # 保存
        output_dir.mkdir(parents=True, exist_ok=True)
        out = output_dir / "budget_template.xlsx"
        wb.save(out)
        return str(out)
    except ImportError:
        return None


def generate_business_ppt_template(output_dir: Path):
    """生成商务PPT模板"""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        from pptx.dml.color import RGBColor
        
        prs = Presentation()
        
        # 第1页：封面
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        if slide.shapes.title:
            slide.shapes.title.text = "项目汇报"
            slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(36)
            slide.shapes.title.text_frame.paragraphs[0].font.bold = True
        
        # 第2页：目录
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        if slide.shapes.title:
            slide.shapes.title.text = "目录"
            slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(28)
        body_shape = None
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 1:
                body_shape = shape
                break
        if body_shape:
            tf = body_shape.text_frame
            tf.text = "项目背景"
            p = tf.add_paragraph()
            p.text = "解决方案"
            p = tf.add_paragraph()
            p.text = "实施计划"
            p = tf.add_paragraph()
            p.text = "预期效果"
            for para in tf.paragraphs:
                para.font.size = Pt(18)
        
        # 第3页：项目背景
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        if slide.shapes.title:
            slide.shapes.title.text = "项目背景"
        body_shape = None
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 1:
                body_shape = shape
                break
        if body_shape:
            tf = body_shape.text_frame
            tf.text = "当前面临的问题和挑战"
            p = tf.add_paragraph()
            p.text = "市场环境分析"
            p = tf.add_paragraph()
            p.text = "竞争对手情况"
            for para in tf.paragraphs:
                para.font.size = Pt(18)
        
        # 第4页：解决方案
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        if slide.shapes.title:
            slide.shapes.title.text = "解决方案"
        body_shape = None
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 1:
                body_shape = shape
                break
        if body_shape:
            tf = body_shape.text_frame
            tf.text = "核心思路"
            p = tf.add_paragraph()
            p.text = "关键措施"
            p = tf.add_paragraph()
            p.text = "创新亮点"
            for para in tf.paragraphs:
                para.font.size = Pt(18)
        
        # 第5页：实施计划
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        if slide.shapes.title:
            slide.shapes.title.text = "实施计划"
        body_shape = None
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 1:
                body_shape = shape
                break
        if body_shape:
            tf = body_shape.text_frame
            tf.text = "第一阶段（1-2周）：准备"
            p = tf.add_paragraph()
            p.text = "第二阶段（3-6周）：执行"
            p = tf.add_paragraph()
            p.text = "第三阶段（7-8周）：验收"
            for para in tf.paragraphs:
                para.font.size = Pt(18)
        
        # 第6页：预期效果
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        if slide.shapes.title:
            slide.shapes.title.text = "预期效果"
        body_shape = None
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 1:
                body_shape = shape
                break
        if body_shape:
            tf = body_shape.text_frame
            tf.text = "效率提升 30%"
            p = tf.add_paragraph()
            p.text = "成本降低 20%"
            p = tf.add_paragraph()
            p.text = "客户满意度显著提高"
            for para in tf.paragraphs:
                para.font.size = Pt(18)
        
        # 第7页：谢谢
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        if slide.shapes.title:
            slide.shapes.title.text = "谢谢"
            slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(40)
            slide.shapes.title.text_frame.paragraphs[0].font.bold = True
            slide.shapes.title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # 保存
        output_dir.mkdir(parents=True, exist_ok=True)
        out = output_dir / "business_ppt_template.pptx"
        prs.save(out)
        return str(out)
    except ImportError:
        return None


def main():
    import argparse
    parser = argparse.ArgumentParser(description="生成 WPS Office 模板文件")
    parser.add_argument("--dir", default=".", help="输出目录（默认当前目录）")
    args = parser.parse_args()
    
    output_dir = Path(args.dir).resolve()
    output_dir.mkdir(parents=True, exist_ok=True)
    
    results = []
    
    # 生成报告模板
    r = generate_report_template(output_dir)
    if r:
        results.append(f"✅ 工作报告模板: {r}")
    else:
        results.append("❌ 工作报告模板: 需要 python-docx（pip install python-docx）")
    
    # 生成预算模板
    r = generate_budget_template(output_dir)
    if r:
        results.append(f"✅ 月度预算模板: {r}")
    else:
        results.append("❌ 月度预算模板: 需要 openpyxl（pip install openpyxl）")
    
    # 生成PPT模板
    r = generate_business_ppt_template(output_dir)
    if r:
        results.append(f"✅ 商务PPT模板: {r}")
    else:
        results.append("❌ 商务PPT模板: 需要 python-pptx（pip install python-pptx）")
    
    print("\n".join(results))


if __name__ == "__main__":
    main()
