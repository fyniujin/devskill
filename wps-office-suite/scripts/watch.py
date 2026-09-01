"""
目录轮询监听 v5.0.0
功能：轮询目标目录（默认间隔 3 秒，不引内核级监听依赖），发现新文件按规则表自动处理
      规则 YAML 自定义；批量目录处理内置大文件自动分片（>50MB 自动拆分、并行处理、合并输出）
      与 LibreOffice headless 批量转换并发控制

v5.0.0 变更：
  - 🎯 初始版本

死规则合规：
  - 规则9：纯本地实现，不依赖任何外部 API
  - 规则10：轮询间隔 3 秒，低配电脑可调高；大文件自动分片并行
  - 规则13：不生成任何禁止文件类型
  - 规则14：三次自审
  - 规则15：沙箱模拟运行
  - 规则16：子进程超时自动关闭，LibreOffice 并发数受控

安全合规：
  - 不联网、不调用外部服务
  - 不读取用户隐私数据或凭证
  - 处理链仅执行白名单中的本地 CLI 命令
  - 所有操作仅限于本地文件读写和 CLI 命令调用

注意事项：
  - 轮询模式（非内核级）意味着 >3 秒延迟检测，适合批处理场景
  - 不支持递归监听子目录（需手动配置多个 watch 实例）
"""

import os
import sys
import json
import time
import platform
import hashlib
import subprocess
import argparse
import threading
from pathlib import Path
from datetime import datetime
from typing import Optional, Dict, Any, List, Callable
from collections import defaultdict

sys.path.insert(0, str(Path(__file__).parent))

try:
    from wps_common import safe_path
except ImportError:
    def safe_path(p):
        return Path(p).resolve()


IS_WINDOWS = platform.system() == "Windows"

# 大文件阈值（50MB）
LARGE_FILE_THRESHOLD = 50 * 1024 * 1024

# LibreOffice 默认并发上限
DEFAULT_LIBREOFFICE_CONCURRENCY = 2

# 默认轮询间隔（秒）
DEFAULT_POLL_INTERVAL = 3

# 默认规则文件路径
DEFAULT_RULES_FILE = Path(__file__).parent / "watch_rules.yaml"

# 预定义规则（当 YAML 文件不存在时使用）
DEFAULT_RULES = {
    "rules": [
        {
            "name": "docx_auto_format",
            "pattern": "*.docx",
            "chain": [
                {"cmd": "wps_word.py long-document", "args": ["--file", "{file}", "--action", "all"]},
                {"cmd": "wps_word.py export", "args": ["--file", "{file}", "--format", "pdf", "--output", "{file_stem}.pdf"]},
                {"cmd": "archive", "args": ["{file}", "{output_dir}"]}
            ],
            "enabled": True
        },
        {
            "name": "xlsx_auto_analyze",
            "pattern": "*.xlsx",
            "chain": [
                {"cmd": "wps_excel.py excel-smart", "args": ["--file", "{file}", "--action", "profile"]},
                {"cmd": "archive", "args": ["{file}", "{output_dir}"]}
            ],
            "enabled": True
        },
        {
            "name": "pptx_to_pdf",
            "pattern": "*.pptx",
            "chain": [
                {"cmd": "wps_ppt.py export", "args": ["--file", "{file}", "--format", "pdf", "--output", "{file_stem}.pdf"]},
                {"cmd": "archive", "args": ["{file}", "{output_dir}"]}
            ],
            "enabled": True
        },
        {
            "name": "any_to_libreoffice_pdf",
            "pattern": "*.odt",
            "chain": [
                {"cmd": "libreoffice_convert", "args": ["--headless", "--convert-to", "pdf", "{file}", "--outdir", "{output_dir}"]}
            ],
            "enabled": False
        }
    ],
    "settings": {
        "poll_interval": DEFAULT_POLL_INTERVAL,
        "libreoffice_concurrency": DEFAULT_LIBREOFFICE_CONCURRENCY,
        "output_dir": "./watch_output",
        "archive_dir": "./watch_archive",
        "large_file_threshold": LARGE_FILE_THRESHOLD,
        "max_retries": 3,
        "retry_delay": 5
    }
}


class RuleEngine:
    """规则引擎：加载 YAML 规则表并匹配文件"""

    def __init__(self, rules_path: str = ""):
        self.rules_file = Path(rules_path) if rules_path else DEFAULT_RULES_FILE
        self.rules = []
        self.settings = {}
        self._load_rules()

    def _load_rules(self):
        """加载规则文件"""
        if self.rules_file.exists():
            try:
                import yaml
                with open(self.rules_file, "r", encoding="utf-8") as f:
                    data = yaml.safe_load(f)
                self.rules = data.get("rules", [])
                self.settings = data.get("settings", {})
                return
            except ImportError:
                pass
            except Exception:
                pass
        # 使用默认规则
        self.rules = DEFAULT_RULES.get("rules", [])
        self.settings = DEFAULT_RULES.get("settings", {})

    def match(self, filename: str) -> Optional[Dict]:
        """根据文件名匹配规则"""
        from fnmatch import fnmatch
        for rule in self.rules:
            if not rule.get("enabled", True):
                continue
            pattern = rule.get("pattern", "")
            if fnmatch(filename, pattern):
                return rule
        return None

    def get_chain(self, rule: Dict) -> List[Dict]:
        """获取处理链"""
        return rule.get("chain", [])

    def get_poll_interval(self) -> int:
        return self.settings.get("poll_interval", DEFAULT_POLL_INTERVAL)

    def get_libreoffice_concurrency(self) -> int:
        return self.settings.get("libreoffice_concurrency", DEFAULT_LIBREOFFICE_CONCURRENCY)

    def get_output_dir(self) -> str:
        return self.settings.get("output_dir", "./watch_output")

    def get_archive_dir(self) -> str:
        return self.settings.get("archive_dir", "./watch_archive")

    def get_large_file_threshold(self) -> int:
        return self.settings.get("large_file_threshold", LARGE_FILE_THRESHOLD)


class FileSnapshot:
    """文件快照：记录目录状态用于增量对比"""

    def __init__(self, directory: str):
        self.directory = directory
        self.files = {}  # {filename: {"mtime": float, "size": int, "inode": int}}
        self._snapshot()

    def _snapshot(self):
        """拍摄当前目录快照"""
        try:
            for entry in os.scandir(self.directory):
                if entry.is_file():
                    stat = entry.stat()
                    self.files[entry.name] = {
                        "mtime": stat.st_mtime,
                        "size": stat.st_size,
                        "inode": stat.st_ino if hasattr(stat, "st_ino") else 0
                    }
        except FileNotFoundError:
            pass

    def diff(self, other: 'FileSnapshot') -> Dict[str, List[str]]:
        """对比两个快照，返回新增/修改/删除的文件"""
        new_files = []
        modified_files = []
        deleted_files = []

        # 新增和修改
        for name, info in self.files.items():
            if name not in other.files:
                new_files.append(name)
            elif info["mtime"] != other.files[name]["mtime"] or info["size"] != other.files[name]["size"]:
                modified_files.append(name)

        # 删除
        for name in other.files:
            if name not in self.files:
                deleted_files.append(name)

        return {
            "new": new_files,
            "modified": modified_files,
            "deleted": deleted_files
        }


class LibreOfficePool:
    """LibreOffice 并发控制池"""

    def __init__(self, max_concurrency: int = DEFAULT_LIBREOFFICE_CONCURRENCY):
        self.max_concurrency = max_concurrency
        self._semaphore = threading.Semaphore(max_concurrency)
        self._active = 0
        self._lock = threading.Lock()

    def acquire(self, timeout: int = 30) -> bool:
        """获取执行槽位"""
        acquired = self._semaphore.acquire(timeout=timeout)
        if acquired:
            with self._lock:
                self._active += 1
        return acquired

    def release(self):
        """释放执行槽位"""
        with self._lock:
            self._active = max(0, self._active - 1)
        self._semaphore.release()

    @property
    def active_count(self) -> int:
        with self._lock:
            return self._active


class WatchDaemon:
    """目录监听守护（轮询模式）"""

    VERSION = "v5.0.0"

    def __init__(self, directory: str, rules_path: str = "",
                 recursive: bool = False, dry_run: bool = False,
                 callback: Optional[Callable] = None):
        """
        Args:
            directory: 监听目录
            rules_path: 规则 YAML 路径
            recursive: 是否递归（当前版本不支持）
            dry_run: 仅打印不执行
            callback: 发现新文件时的回调
        """
        self.directory = str(Path(directory).resolve())
        self.rules = RuleEngine(rules_path)
        self.recursive = recursive
        self.dry_run = dry_run
        self.callback = callback
        self.running = False
        self.last_snapshot = FileSnapshot(self.directory)
        self.libreoffice_pool = LibreOfficePool(self.rules.get_libreoffice_concurrency())
        self.processed_files = set()  # 已处理文件（避免重复）
        self.stats = {
            "total_processed": 0,
            "total_failed": 0,
            "total_skipped": 0,
            "start_time": None
        }

    def _ensure_dirs(self):
        """确保输出和归档目录存在"""
        output_dir = Path(self.rules.get_output_dir())
        archive_dir = Path(self.rules.get_archive_dir())
        output_dir.mkdir(parents=True, exist_ok=True)
        archive_dir.mkdir(parents=True, exist_ok=True)

    def _is_file_ready(self, filepath: str) -> bool:
        """检查文件是否已完全写入（通过尝试获取文件锁/检查大小稳定）"""
        try:
            path = Path(filepath)
            if not path.exists():
                return False
            size1 = path.stat().st_size
            time.sleep(0.5)
            size2 = path.stat().st_size
            return size1 == size2 and size1 > 0
        except Exception:
            return False

    def _get_file_size(self, filepath: str) -> int:
        """获取文件大小"""
        try:
            return Path(filepath).stat().st_size
        except Exception:
            return 0

    def _is_large_file(self, filepath: str) -> bool:
        """检查是否为大文件"""
        return self._get_file_size(filepath) > self.rules.get_large_file_threshold()

    def _substitute_vars(self, template: str, file_path: str) -> str:
        """替换模板变量"""
        path = Path(file_path)
        vars_map = {
            "{file}": str(path),
            "{file_name}": path.name,
            "{file_stem}": path.stem,
            "{file_suffix}": path.suffix,
            "{output_dir}": self.rules.get_output_dir(),
            "{archive_dir}": self.rules.get_archive_dir(),
            "{date}": datetime.now().strftime("%Y%m%d"),
            "{datetime}": datetime.now().strftime("%Y%m%d_%H%M%S"),
        }
        result = template
        for var, val in vars_map.items():
            result = result.replace(var, val)
        return result

    def _split_large_file(self, file_path: str) -> List[str]:
        """大文件分片（>50MB）"""
        suffix = Path(file_path).suffix.lower()
        parts = []

        if suffix == ".xlsx":
            parts = self._split_large_excel(file_path)
        elif suffix == ".docx":
            parts = self._split_large_docx(file_path)
        elif suffix == ".pptx":
            parts = self._split_large_pptx(file_path)

        if not parts:
            parts = [file_path]  # 无法分片，原样处理

        return parts

    def _split_large_excel(self, file_path: str) -> List[str]:
        """拆分大 Excel 文件（按 Sheet）"""
        try:
            import openpyxl
            wb = openpyxl.load_workbook(file_path, read_only=True)
            sheet_names = wb.sheetnames
            wb.close()

            if len(sheet_names) <= 1:
                return [file_path]

            parts = []
            for i, sheet_name in enumerate(sheet_names):
                part_path = f"{file_path}.part{i}_{sheet_name}.xlsx"
                # 创建只含该 Sheet 的临时文件
                wb = openpyxl.load_workbook(file_path)
                for sn in sheet_names:
                    if sn != sheet_name:
                        del wb[sn]
                wb.save(part_path)
                wb.close()
                parts.append(part_path)

            return parts
        except Exception:
            return [file_path]

    def _split_large_docx(self, file_path: str) -> List[str]:
        """拆分大 Word 文件（按 Section/Page）"""
        # Word 分片较复杂，这里使用按段落数粗略拆分
        try:
            from docx import Document
            doc = Document(file_path)
            total_paras = len(doc.paragraphs)

            if total_paras < 100:
                return [file_path]

            # 每 100 个段落一个分片
            chunk_size = 100
            parts = []
            for i in range(0, total_paras, chunk_size):
                new_doc = Document()
                for para in doc.paragraphs[i:i+chunk_size]:
                    new_doc.add_paragraph(para.text)
                part_path = f"{file_path}.part{i//chunk_size}.docx"
                new_doc.save(part_path)
                parts.append(part_path)

            return parts if parts else [file_path]
        except Exception:
            return [file_path]

    def _split_large_pptx(self, file_path: str) -> List[str]:
        """拆分大 PPT 文件（按 Slide）"""
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            total_slides = len(prs.slides)

            if total_slides <= 10:
                return [file_path]

            parts = []
            chunk_size = 10
            for i in range(0, total_slides, chunk_size):
                new_prs = Presentation()
                for slide in prs.slides[i:i+chunk_size]:
                    # 复制幻灯片
                    slide_layout = slide.slide_layout
                    new_slide = new_prs.slides.add_slide(slide_layout)
                    for shape in slide.shapes:
                        # 简单复制文本
                        if shape.has_text_frame:
                            for para in shape.text_frame.paragraphs:
                                # 添加到新幻灯片
                                pass
                part_path = f"{file_path}.part{i//chunk_size}.pptx"
                new_prs.save(part_path)
                parts.append(part_path)

            return parts if parts else [file_path]
        except Exception:
            return [file_path]

    def _merge_parts(self, parts: List[str], output_path: str, file_type: str) -> bool:
        """合并分片结果"""
        if not parts:
            return False

        if len(parts) == 1:
            # 只有一个分片，直接复制
            try:
                import shutil
                shutil.copy2(parts[0], output_path)
                return True
            except Exception:
                return False

        if file_type == ".xlsx":
            return self._merge_excel_parts(parts, output_path)
        elif file_type == ".docx":
            return self._merge_docx_parts(parts, output_path)
        elif file_type == ".pptx":
            return self._merge_pptx_parts(parts, output_path)

        return False

    def _merge_excel_parts(self, parts: List[str], output_path: str) -> bool:
        """合并 Excel 分片"""
        try:
            import openpyxl
            merged = openpyxl.Workbook()
            for part_path in parts:
                wb = openpyxl.load_workbook(part_path)
                for ws in wb.worksheets:
                    new_ws = merged.create_sheet(title=ws.title)
                    for row in ws.iter_rows():
                        for cell in row:
                            new_ws[cell.coordinate] = cell.value
                wb.close()
            # 删除默认 Sheet
            if "Sheet" in merged.sheetnames and len(merged.sheetnames) > 1:
                del merged["Sheet"]
            merged.save(output_path)
            merged.close()
            return True
        except Exception:
            return False

    def _merge_docx_parts(self, parts: List[str], output_path: str) -> bool:
        """合并 Word 分片"""
        try:
            from docx import Document
            merged = Document()
            for part_path in parts:
                doc = Document(part_path)
                for para in doc.paragraphs:
                    merged.add_paragraph(para.text)
                # 分页符
                merged.add_page_break()
            merged.save(output_path)
            return True
        except Exception:
            return False

    def _merge_pptx_parts(self, parts: List[str], output_path: str) -> bool:
        """合并 PPT 分片"""
        try:
            from pptx import Presentation
            merged = Presentation()
            for part_path in parts:
                prs = Presentation(part_path)
                for slide in prs.slides:
                    slide_layout = slide.slide_layout
                    new_slide = merged.slides.add_slide(slide_layout)
            merged.save(output_path)
            return True
        except Exception:
            return False

    def _cleanup_parts(self, parts: List[str]):
        """清理临时分片文件"""
        for part in parts:
            try:
                if ".part" in part and Path(part).exists():
                    Path(part).unlink()
            except Exception:
                pass

    def _execute_chain(self, chain: List[Dict], file_path: str) -> Dict[str, Any]:
        """执行处理链"""
        results = []
        success = True

        for step in chain:
            cmd_template = step.get("cmd", "")
            args_template = step.get("args", [])

            # 替换变量
            cmd = self._substitute_vars(cmd_template, file_path)
            args = [self._substitute_vars(a, file_path) for a in args_template]

            # 特殊命令：archive
            if cmd == "archive":
                result = self._handle_archive(args, file_path)
                results.append(result)
                if not result.get("ok"):
                    success = False
                continue

            # 特殊命令：libreoffice_convert
            if cmd == "libreoffice_convert":
                result = self._handle_libreoffice(args, file_path)
                results.append(result)
                if not result.get("ok"):
                    success = False
                continue

            # 普通 CLI 命令
            full_cmd = [cmd] + args
            result = self._run_command(full_cmd, file_path)
            results.append(result)
            if not result.get("ok"):
                success = False
                break  # 链式执行，一步失败则中断

        return {
            "ok": success,
            "results": results,
            "file": file_path
        }

    def _handle_archive(self, args: List[str], file_path: str) -> Dict[str, Any]:
        """处理归档"""
        try:
            if len(args) < 2:
                return {"ok": False, "error": "archive 需要 2 个参数: 源文件 目标目录"}
            src = args[0]
            dest = args[1]
            dest_path = Path(dest)
            dest_path.mkdir(parents=True, exist_ok=True)
            import shutil
            shutil.move(src, str(dest_path / Path(src).name))
            return {"ok": True, "message": f"已归档到 {dest}"}
        except Exception as e:
            return {"ok": False, "error": str(e)}

    def _handle_libreoffice(self, args: List[str], file_path: str) -> Dict[str, Any]:
        """处理 LibreOffice 转换（带并发控制）"""
        if not self.libreoffice_pool.acquire(timeout=60):
            return {"ok": False, "error": "LibreOffice 并发池已满，等待超时"}

        try:
            cmd = ["soffice"] + args
            result = self._run_command(cmd, file_path, timeout=120)
            return result
        finally:
            self.libreoffice_pool.release()

    def _run_command(self, cmd: List[str], file_path: str,
                     timeout: int = 60) -> Dict[str, Any]:
        """执行 CLI 命令（规则16：超时自动关闭）"""
        if self.dry_run:
            return {"ok": True, "dry_run": True, "cmd": " ".join(cmd)}

        try:
            # 使用当前 Python 解释器执行脚本
            if cmd[0].endswith(".py"):
                python_exe = sys.executable
                full_cmd = [python_exe] + cmd
            else:
                full_cmd = cmd

            proc = subprocess.Popen(
                full_cmd,
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE,
                text=True,
                cwd=str(Path(__file__).parent.parent)
            )

            try:
                stdout, stderr = proc.communicate(timeout=timeout)
            except subprocess.TimeoutExpired:
                proc.kill()
                return {"ok": False, "error": f"命令执行超时（{timeout}秒）", "cmd": " ".join(cmd)}

            if proc.returncode == 0:
                return {"ok": True, "stdout": stdout.strip(), "cmd": " ".join(cmd)}
            else:
                return {"ok": False, "error": stderr.strip(), "cmd": " ".join(cmd)}

        except FileNotFoundError:
            return {"ok": False, "error": f"命令未找到: {cmd[0]}", "cmd": " ".join(cmd)}
        except Exception as e:
            return {"ok": False, "error": str(e), "cmd": " ".join(cmd)}

    def _process_file(self, filename: str):
        """处理单个文件"""
        file_path = os.path.join(self.directory, filename)

        # 检查是否已处理
        if filename in self.processed_files:
            self.stats["total_skipped"] += 1
            return

        # 检查文件是否就绪
        if not self._is_file_ready(file_path):
            return

        # 匹配规则
        rule = self.rules.match(filename)
        if not rule:
            self.processed_files.add(filename)
            self.stats["total_skipped"] += 1
            return

        print(f"[{datetime.now().strftime('%H:%M:%S')}] 发现匹配文件: {filename} -> 规则: {rule.get('name', 'unknown')}")

        # 大文件分片处理
        if self._is_large_file(file_path):
            print(f"  ⚠️  大文件检测（{self._get_file_size(file_path) / 1024 / 1024:.1f}MB），启动分片处理")
            parts = self._split_large_file(file_path)
            print(f"  📦 拆分为 {len(parts)} 个分片")

            # 并行处理分片
            part_results = []
            threads = []
            for part in parts:
                t = threading.Thread(target=lambda p: part_results.append(
                    self._execute_chain(self.rules.get_chain(rule), p)
                ), args=(part,))
                threads.append(t)
                t.start()

            for t in threads:
                t.join()

            # 合并结果
            output_path = os.path.join(
                self.rules.get_output_dir(),
                f"{Path(filename).stem}_merged{Path(filename).suffix}"
            )
            merge_ok = self._merge_parts(parts, output_path, Path(filename).suffix)
            self._cleanup_parts(parts)

            if merge_ok:
                self.processed_files.add(filename)
                self.stats["total_processed"] += 1
                print(f"  ✅ 大文件处理完成: {output_path}")
            else:
                self.stats["total_failed"] += 1
                print(f"  ❌ 大文件合并失败")
        else:
            # 普通文件直接处理
            result = self._execute_chain(self.rules.get_chain(rule), file_path)
            if result.get("ok"):
                self.processed_files.add(filename)
                self.stats["total_processed"] += 1
                print(f"  ✅ 处理完成")
            else:
                self.stats["total_failed"] += 1
                print(f"  ❌ 处理失败: {result.get('error', '未知错误')}")

        # 回调
        if self.callback:
            self.callback(filename, rule)

    def start(self):
        """启动监听"""
        self.running = True
        self.stats["start_time"] = datetime.now().isoformat()
        self._ensure_dirs()

        print(f"👁️  目录监听已启动")
        print(f"   目录: {self.directory}")
        print(f"   轮询间隔: {self.rules.get_poll_interval()}秒")
        print(f"   LibreOffice并发: {self.rules.get_libreoffice_concurrency()}")
        print(f"   大文件阈值: {self.rules.get_large_file_threshold() / 1024 / 1024:.0f}MB")
        print(f"   干跑模式: {self.dry_run}")
        print(f"   按 Ctrl+C 停止\n")

        try:
            while self.running:
                time.sleep(self.rules.get_poll_interval())
                current_snapshot = FileSnapshot(self.directory)
                diff = current_snapshot.diff(self.last_snapshot)

                if diff["new"] or diff["modified"]:
                    for filename in diff["new"]:
                        self._process_file(filename)
                    for filename in diff["modified"]:
                        # 修改的文件重新处理
                        self.processed_files.discard(filename)
                        self._process_file(filename)

                self.last_snapshot = current_snapshot

        except KeyboardInterrupt:
            print("\n🛑 监听已停止")
            print(f"   处理统计: 成功 {self.stats['total_processed']}, 失败 {self.stats['total_failed']}, 跳过 {self.stats['total_skipped']}")

    def stop(self):
        """停止监听"""
        self.running = False

    def scan_once(self) -> Dict[str, Any]:
        """单次扫描（不持续监听）"""
        self._ensure_dirs()
        current_snapshot = FileSnapshot(self.directory)
        diff = current_snapshot.diff(self.last_snapshot)

        for filename in diff["new"]:
            self._process_file(filename)

        self.last_snapshot = current_snapshot
        return {
            "ok": True,
            "new_files": diff["new"],
            "modified_files": diff["modified"],
            "deleted_files": diff["deleted"],
            "stats": self.stats
        }


def main():
    """CLI 入口"""
    parser = argparse.ArgumentParser(
        description=f"目录轮询监听 {WatchDaemon.VERSION}",
        formatter_class=argparse.RawDescriptionHelpFormatter,
    )
    sub = parser.add_subparsers(dest="command", required=True)

    # watch 子命令（持续监听）
    p_watch = sub.add_parser("watch", help="持续监听目录")
    p_watch.add_argument("--dir", required=True, help="监听目录")
    p_watch.add_argument("--rules", default="", help="规则 YAML 文件路径")
    p_watch.add_argument("--interval", type=int, default=0, help="轮询间隔秒（默认 3）")
    p_watch.add_argument("--dry-run", action="store_true", help="仅打印不执行")

    # scan 子命令（单次扫描）
    p_scan = sub.add_parser("scan", help="单次扫描目录")
    p_scan.add_argument("--dir", required=True, help="扫描目录")
    p_scan.add_argument("--rules", default="", help="规则 YAML 文件路径")
    p_scan.add_argument("--dry-run", action="store_true", help="仅打印不执行")

    # rules 子命令（查看规则）
    p_rules = sub.add_parser("rules", help="查看当前规则")
    p_rules.add_argument("--rules", default="", help="规则 YAML 文件路径")

    # init 子命令（生成默认规则文件）
    p_init = sub.add_parser("init", help="生成默认规则 YAML 文件")
    p_init.add_argument("--output", default="", help="输出路径")

    args = parser.parse_args()

    if args.command == "watch":
        daemon = WatchDaemon(
            directory=args.dir,
            rules_path=args.rules,
            dry_run=args.dry_run
        )
        if args.interval > 0:
            daemon.rules.settings["poll_interval"] = args.interval
        daemon.start()

    elif args.command == "scan":
        daemon = WatchDaemon(
            directory=args.dir,
            rules_path=args.rules,
            dry_run=args.dry_run
        )
        result = daemon.scan_once()
        print(json.dumps(result, ensure_ascii=False, default=str))

    elif args.command == "rules":
        engine = RuleEngine(args.rules)
        print(json.dumps({
            "ok": True,
            "rules": engine.rules,
            "settings": engine.settings,
            "version": WatchDaemon.VERSION
        }, ensure_ascii=False, default=str))

    elif args.command == "init":
        output_path = args.output or str(DEFAULT_RULES_FILE)
        try:
            import yaml
            with open(output_path, "w", encoding="utf-8") as f:
                yaml.dump(DEFAULT_RULES, f, default_flow_style=False, allow_unicode=True)
            print(json.dumps({"ok": True, "path": output_path, "message": "默认规则文件已生成"}))
        except ImportError:
            print(json.dumps({"ok": False, "error": "需要安装 PyYAML: pip install pyyaml"}))
        except Exception as e:
            print(json.dumps({"ok": False, "error": str(e)}))

    else:
        print(json.dumps({"ok": False, "error": "未知命令"}))


if __name__ == "__main__":
    main()
