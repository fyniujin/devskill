"""
Markdown → Word/PPT 转换器 v4.7.0
功能：Markdown 文件转换为 Word 文档或 PPT 演示

v4.7.0 变更:
  - 🎯 Markdown → Word（保留标题层级、列表、表格、代码块）
  - 🎯 Markdown → PPT（## 标题分页、内容自动拆分）
  - 🎯 纯本地实现，不读取外部凭证或 API Key
  - 🎯 批量目录转换
  - 🎯 硬件自适应（低配减少并发）
"""

import os
import sys
import json
import re
from pathlib import Path
from typing import Dict, List, Optional, Any

sys.path.insert(0, str(Path(__file__).parent))

try:
    from wps_common import safe_path, get_hardware_info
except ImportError:
    def safe_path(p): return Path(p)
    def get_hardware_info(): return {"cpu_cores": 4, "memory_gb": 8, "level": "medium"}

try:
    import docx
    from docx.shared import Pt, Inches
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    import pptx
    from pptx.util import Inches as PptxInches, Pt as PptxPt
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False


class MDParser:
    """Markdown 解析器"""
    
    def __init__(self):
        pass
    
    def parse(self, content: str) -> List[Dict]:
        """解析 Markdown 内容为结构化节点"""
        nodes = []
        lines = content.split('\n')
        i = 0
        
        while i < len(lines):
            line = lines[i]
            
            # 空行
            if not line.strip():
                i += 1
                continue
            
            # 标题
            heading_match = re.match(r'^(#{1,6})\s+(.*)', line)
            if heading_match:
                level = len(heading_match.group(1))
                text = heading_match.group(2).strip()
                nodes.append({
                    "type": "heading",
                    "level": level,
                    "text": text,
                })
                i += 1
                continue
            
            # 代码块
            if line.strip().startswith('```'):
                code_lines = []
                i += 1
                while i < len(lines) and not lines[i].strip().startswith('```'):
                    code_lines.append(lines[i])
                    i += 1
                i += 1  # skip closing ```
                nodes.append({
                    "type": "code",
                    "text": '\n'.join(code_lines),
                })
                continue
            
            # 表格
            if '|' in line and i + 1 < len(lines) and re.match(r'^[\s|:-]+$', lines[i+1]):
                # 表头
                headers = [c.strip() for c in line.strip('| ').split('|')]
                i += 2  # skip header and separator
                rows = []
                while i < len(lines) and '|' in lines[i] and lines[i].strip():
                    row = [c.strip() for c in lines[i].strip('| ').split('|')]
                    rows.append(row)
                    i += 1
                nodes.append({
                    "type": "table",
                    "headers": headers,
                    "rows": rows,
                })
                continue
            
            # 无序列表
            if re.match(r'^[\s]*[-*+]\s', line):
                list_items = []
                while i < len(lines) and re.match(r'^[\s]*[-*+]\s', lines[i]):
                    list_items.append(re.sub(r'^[\s]*[-*+]\s+', '', lines[i]))
                    i += 1
                nodes.append({
                    "type": "unordered_list",
                    "items": list_items,
                })
                continue
            
            # 有序列表
            if re.match(r'^[\s]*\d+\.\s', line):
                list_items = []
                while i < len(lines) and re.match(r'^[\s]*\d+\.\s', lines[i]):
                    list_items.append(re.sub(r'^[\s]*\d+\.\s+', '', lines[i]))
                    i += 1
                nodes.append({
                    "type": "ordered_list",
                    "items": list_items,
                })
                continue
            
            # 分割线
            if re.match(r'^[\s]*---[\s]*$', line) or re.match(r'^[\s]*\*\*[\s]*$', line):
                nodes.append({"type": "hr"})
                i += 1
                continue
            
            # 引用块
            if line.strip().startswith('>'):
                quote_lines = []
                while i < len(lines) and lines[i].strip().startswith('>'):
                    quote_lines.append(re.sub(r'^>\s?', '', lines[i]))
                    i += 1
                nodes.append({
                    "type": "quote",
                    "text": '\n'.join(quote_lines),
                })
                continue
            
            # 普通段落（可能有多行）
            para_lines = [line]
            i += 1
            while i < len(lines) and lines[i].strip() and not re.match(r'^(#{1,6}\s|```|\||[-*+\s]|\d+\.\s|---|\*\*\*|>)', lines[i]):
                para_lines.append(lines[i])
                i += 1
            nodes.append({
                "type": "paragraph",
                "text": ' '.join(para_lines),
            })
        
        return nodes


class MDConverter:
    """Markdown → Word/PPT 转换器"""
    
    def __init__(self):
        self.parser = MDParser()
        self.hw = get_hardware_info()
    
    def md_to_docx(self, input_path: str, output_path: str) -> Dict:
        """Markdown → Word"""
        if not HAS_DOCX:
            return {"success": False, "error": "python-docx 未安装"}
        
        try:
            content = Path(input_path).read_text(encoding='utf-8')
            nodes = self.parser.parse(content)
            
            doc = docx.Document()
            
            for node in nodes:
                ntype = node["type"]
                
                if ntype == "heading":
                    level = min(node["level"], 9)
                    heading = doc.add_heading(node["text"], level=level)
                
                elif ntype == "paragraph":
                    if node["text"].strip():
                        para = doc.add_paragraph()
                        self._add_formatted_text(para, node["text"])
                
                elif ntype == "code":
                    para = doc.add_paragraph()
                    run = para.add_run(node["text"])
                    run.font.name = 'Courier New'
                    run.font.size = Pt(10)
                
                elif ntype == "unordered_list":
                    for item in node["items"]:
                        para = doc.add_paragraph(item, style='List Bullet')
                
                elif ntype == "ordered_list":
                    for item in node["items"]:
                        para = doc.add_paragraph(item, style='List Number')
                
                elif ntype == "table":
                    table = doc.add_table(rows=1 + len(node["rows"]), cols=len(node["headers"]))
                    table.style = 'Table Grid'
                    # 表头
                    for j, header in enumerate(node["headers"]):
                        table.rows[0].cells[j].text = header
                    # 数据行
                    for i, row in enumerate(node["rows"]):
                        for j, cell in enumerate(row):
                            if j < len(node["headers"]):
                                table.rows[i+1].cells[j].text = cell
                
                elif ntype == "quote":
                    para = doc.add_paragraph(node["text"])
                    para.paragraph_format.left_indent = Inches(0.5)
                
                elif ntype == "hr":
                    doc.add_paragraph('─' * 40)
            
            doc.save(output_path)
            
            return {
                "success": True,
                "output": output_path,
                "nodes_converted": len(nodes),
            }
        except Exception as e:
            return {"success": False, "error": str(e)}
    
    def md_to_pptx(self, input_path: str, output_path: str) -> Dict:
        """Markdown → PPT"""
        if not HAS_PPTX:
            return {"success": False, "error": "python-pptx 未安装"}
        
        try:
            content = Path(input_path).read_text(encoding='utf-8')
            nodes = self.parser.parse(content)
            
            prs = pptx.Presentation()
            
            # 第一张标题 slide
            slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title Slide
            title_shape = slide.shapes.title
            subtitle_shape = slide.placeholders[1] if len(slide.placeholders) > 1 else None
            
            # 找第一个 # 标题作为标题
            first_heading = None
            for node in nodes:
                if node["type"] == "heading" and node["level"] == 1:
                    first_heading = node["text"]
                    break
            
            title_shape.text = first_heading or Path(input_path).stem
            if subtitle_shape:
                subtitle_shape.text = f"由 {Path(input_path).name} 自动生成"
            
            # 后续 ## 标题分页
            current_slide = None
            content_shapes = []
            
            for node in nodes:
                ntype = node["type"]
                
                if ntype == "heading" and node["level"] == 2:
                    # 新建幻灯片
                    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title + Content
                    slide.shapes.title.text = node["text"]
                    content_shapes = [s for s in slide.shapes if s != slide.shapes.title]
                    current_slide = slide
                
                elif ntype == "heading" and node["level"] == 3:
                    # 三级标题：在当前 slide 中添加
                    if current_slide:
                        if content_shapes:
                            tf = content_shapes[0].text_frame
                            tf.text = node["text"]
                        else:
                            txBox = current_slide.shapes.add_textbox(PptxInches(1), PptxInches(2), PptxInches(8), PptxInches(1))
                            tf = txBox.text_frame
                            tf.text = node["text"]
                
                elif ntype in ("paragraph", "unordered_list", "ordered_list"):
                    text = ""
                    if ntype == "paragraph":
                        text = node["text"]
                    else:
                        text = '\n'.join(f"• {item}" for item in node["items"])
                    
                    if text.strip():
                        if current_slide:
                            if content_shapes:
                                tf = content_shapes[0].text_frame
                                if tf.text:
                                    tf.text = tf.text + '\n' + text
                                else:
                                    tf.text = text
                            else:
                                txBox = current_slide.shapes.add_textbox(PptxInches(1), PptxInches(2), PptxInches(8), PptxInches(5))
                                tf = txBox.text_frame
                                tf.text = text
                                tf.word_wrap = True
                        else:
                            # 没有 ## 标题时，新建一个默认 slide
                            slide = prs.slides.add_slide(prs.slide_layouts[1])
                            slide.shapes.title.text = "内容"
                            current_slide = slide
                            content_shapes = [s for s in slide.shapes if s != slide.shapes.title]
                            if content_shapes:
                                content_shapes[0].text_frame.text = text
            
            prs.save(output_path)
            
            return {
                "success": True,
                "output": output_path,
                "slides": len(prs.slides),
                "nodes_converted": len(nodes),
            }
        except Exception as e:
            return {"success": False, "error": str(e)}
    
    def _add_formatted_text(self, para, text: str):
        """添加带格式的文本（支持粗体、斜体、代码）"""
        # 简单的内联格式解析
        parts = re.split(r'(\*\*.*?\*\*|\*.*?\*|`.*?`|_.*?_)', text)
        for part in parts:
            if not part:
                continue
            if part.startswith('**') and part.endswith('**'):
                run = para.add_run(part[2:-2])
                run.bold = True
            elif part.startswith('*') and part.endswith('*'):
                run = para.add_run(part[1:-1])
                run.italic = True
            elif part.startswith('`') and part.endswith('`'):
                run = para.add_run(part[1:-1])
                run.font.name = 'Courier New'
                run.font.size = Pt(10)
            elif part.startswith('_') and part.endswith('_'):
                run = para.add_run(part[1:-1])
                run.italic = True
            else:
                para.add_run(part)
    
    def batch_convert(self, input_dir: str, output_dir: str,
                     output_format: str = "docx") -> Dict:
        """批量转换目录中的 Markdown 文件"""
        input_dir = Path(input_dir)
        output_dir = Path(output_dir)
        output_dir.mkdir(parents=True, exist_ok=True)
        
        md_files = list(input_dir.glob("*.md")) + list(input_dir.glob("*.markdown"))
        
        results = []
        for md_file in md_files:
            output_path = str(output_dir / f"{md_file.stem}.{output_format}")
            if output_format == "pptx":
                result = self.md_to_pptx(str(md_file), output_path)
            else:
                result = self.md_to_docx(str(md_file), output_path)
            results.append({"file": md_file.name, **result})
        
        success = sum(1 for r in results if r.get("success"))
        return {
            "success": True,
            "total": len(results),
            "success_count": success,
            "results": results,
        }


def _cli():
    """CLI 入口"""
    import argparse
    
    parser = argparse.ArgumentParser(description="Markdown → Word/PPT 转换器 v4.7.0")
    sub = parser.add_subparsers(dest="command", required=True)
    
    # to-docx
    p = sub.add_parser("to-docx", help="MD → Word")
    p.add_argument("--file", required=True)
    p.add_argument("--output", required=True)
    
    # to-pptx
    p = sub.add_parser("to-pptx", help="MD → PPT")
    p.add_argument("--file", required=True)
    p.add_argument("--output", required=True)
    
    # batch
    p = sub.add_parser("batch", help="批量转换")
    p.add_argument("--input-dir", required=True)
    p.add_argument("--output-dir", required=True)
    p.add_argument("--format", default="docx", choices=["docx", "pptx"])
    
    args = parser.parse_args()
    
    converter = MDConverter()
    
    if args.command == "to-docx":
        result = converter.md_to_docx(args.file, args.output)
        print(json.dumps(result, ensure_ascii=False))
    
    elif args.command == "to-pptx":
        result = converter.md_to_pptx(args.file, args.output)
        print(json.dumps(result, ensure_ascii=False))
    
    elif args.command == "batch":
        result = converter.batch_convert(args.input_dir, args.output_dir, args.format)
        print(json.dumps(result, ensure_ascii=False))


if __name__ == "__main__":
    _cli()
