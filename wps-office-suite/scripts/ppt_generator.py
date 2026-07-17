"""
PPT 智能生成器 v4.1.0（重构版）
多源输入 + 智能增强 + 视觉美化 + 排练辅助

不依赖外部 API（规则9：基础功能自研 + 外部API按需接入）
"""
import re
import json
from pathlib import Path
from typing import List, Dict, Optional, Tuple
from datetime import datetime


# ==================== 统一文档模型 ====================

class Section:
    """文档章节节点"""
    def __init__(self, title: str, level: int = 1, content: List[str] = None, data: Dict = None):
        self.title = title
        self.level = level
        self.content = content or []
        self.data = data or {}

    def to_dict(self) -> Dict:
        return {
            "title": self.title,
            "level": self.level,
            "content": self.content,
            "data": self.data,
        }


class Document:
    """统一文档模型（所有输入格式最终转换为此）"""
    def __init__(self, title: str = "", sections: List[Section] = None, metadata: Dict = None):
        self.title = title
        self.sections = sections or []
        self.metadata = metadata or {}

    def to_dict(self) -> Dict:
        return {
            "title": self.title,
            "sections": [s.to_dict() for s in self.sections],
            "metadata": self.metadata,
        }


# ==================== 输入适配层 ====================

class WordAdapter:
    """Word 文档适配器（python-docx）"""

    @staticmethod
    def parse(filepath: str) -> Document:
        try:
            from docx import Document as DocxDocument
            doc = DocxDocument(filepath)
            document = Document(title=Path(filepath).stem)

            current_section = None
            for para in doc.paragraphs:
                text = para.text.strip()
                if not text:
                    continue

                style_name = para.style.name if para.style else ""
                is_heading = style_name.startswith("Heading") or style_name.startswith("标题")

                if is_heading:
                    match = re.search(r"(\d+)", style_name)
                    level = int(match.group(1)) if match else 1
                    if current_section:
                        document.sections.append(current_section)
                    current_section = Section(title=text, level=level)
                elif current_section:
                    current_section.content.append(text)

            if current_section:
                document.sections.append(current_section)

            return document
        except ImportError:
            return Document(title="Error", metadata={"error": "需要安装 python-docx"})
        except Exception as e:
            return Document(title="Error", metadata={"error": str(e)})


class MarkdownAdapter:
    """Markdown 文档适配器"""

    @staticmethod
    def parse(filepath: str) -> Document:
        try:
            text = Path(filepath).read_text(encoding="utf-8")
            return MarkdownAdapter.parse_text(text, title=Path(filepath).stem)
        except Exception as e:
            return Document(title="Error", metadata={"error": str(e)})

    @staticmethod
    def parse_text(text: str, title: str = "") -> Document:
        lines = text.strip().split("\n")
        document = Document(title=title)
        current_section = None

        for line in lines:
            line = line.strip()
            if not line:
                continue

            # 匹配 Markdown 标题
            heading_match = re.match(r"^(#{1,6})\s+(.+)$", line)
            if heading_match:
                level = len(heading_match.group(1))
                heading_text = heading_match.group(2).strip()
                if current_section:
                    document.sections.append(current_section)
                current_section = Section(title=heading_text, level=level)
            elif current_section:
                # 移除 Markdown 格式符号
                clean_line = re.sub(r"[*_`~\[\]]", "", line).strip()
                if clean_line:
                    current_section.content.append(clean_line)

        if current_section:
            document.sections.append(current_section)

        return document


class MindMapAdapter:
    """思维导图适配器（支持 .mm XML 格式和 .xmind）"""

    @staticmethod
    def parse(filepath: str) -> Document:
        try:
            content = Path(filepath).read_text(encoding="utf-8")
            # 简单 XML 解析（提取 topic 标签）
            topics = re.findall(r'<topic[^>]*>.*?<title>(.*?)</title>.*?</topic>', content, re.DOTALL)
            if not topics:
                # 尝试 xmind 格式
                topics = re.findall(r'text="(.*?)"', content)

            document = Document(title=Path(filepath).stem)
            for i, topic in enumerate(topics):
                topic = topic.strip()
                if topic:
                    document.sections.append(Section(title=topic, level=1, content=[]))

            return document
        except Exception as e:
            return Document(title="Error", metadata={"error": str(e)})


class TextAdapter:
    """纯文本大纲适配器"""

    @staticmethod
    def parse(filepath: str) -> Document:
        try:
            text = Path(filepath).read_text(encoding="utf-8")
            return TextAdapter.parse_text(text, title=Path(filepath).stem)
        except Exception as e:
            return Document(title="Error", metadata={"error": str(e)})

    @staticmethod
    def parse_text(text: str, title: str = "") -> Document:
        lines = text.strip().split("\n")
        document = Document(title=title)
        current_section = None

        for line in lines:
            line = line.strip()
            if not line:
                continue

            # 检测标题特征
            is_heading = False
            level = 1

            # 数字编号：1. 或 1、或 第一章
            if re.match(r"^第[一二三四五六七八九十百千]+[章节部分]", line):
                is_heading = True
                level = 1
            elif re.match(r"^\d+[\.\、]", line):
                is_heading = True
                level = 1
            elif re.match(r"^[\(（]\d+[\)）]", line):
                is_heading = True
                level = 2
            elif line.startswith("#") or line.startswith("【"):
                is_heading = True
                level = 1
            elif len(line) < 30 and not line.endswith(("，", "。", "？", "！")):
                is_heading = True
                level = 2

            if is_heading:
                if current_section:
                    document.sections.append(current_section)
                current_section = Section(title=line, level=level)
            elif current_section:
                current_section.content.append(line)

        if current_section:
            document.sections.append(current_section)

        return document


# ==================== 输入适配工厂 ====================

def parse_input(input_path: str, input_type: str = "auto") -> Document:
    """
    通用输入解析工厂

    Args:
        input_path: 输入文件路径
        input_type: 输入类型（auto/word/markdown/mindmap/text）
    """
    if input_type == "auto":
        ext = Path(input_path).suffix.lower()
        type_map = {
            ".docx": "word",
            ".md": "markdown",
            ".markdown": "markdown",
            ".mm": "mindmap",
            ".xmind": "mindmap",
            ".txt": "text",
        }
        input_type = type_map.get(ext, "text")

    adapters = {
        "word": WordAdapter,
        "markdown": MarkdownAdapter,
        "mindmap": MindMapAdapter,
        "text": TextAdapter,
    }

    adapter = adapters.get(input_type, TextAdapter)
    return adapter.parse(input_path)


# ==================== 演讲者备注生成器 ====================

class SpeakerNotesGenerator:
    """
    演讲者备注生成器（双模式）
    - 模板引擎（默认，规则9合规）
    - 外部 LLM（按需接入，用户配置后启用）
    """

    # 开场白模板
    OPENING_TEMPLATES = [
        "各位好，今天我将为大家分享《{title}》。",
        "大家好，接下来由我为大家讲解《{title}》。",
        "欢迎各位，今天的分享主题是《{title}》。",
    ]

    # 过渡语模板
    TRANSITION_TEMPLATES = [
        "接下来，我们来看{section}。",
        "下面进入{section}部分。",
        "让我们继续探讨{section}。",
        "现在转向{section}。",
    ]

    # 要点拓展模板
    EXPANSION_TEMPLATES = [
        "关于{point}，需要特别注意的是……",
        "在{point}方面，我们可以从以下几个角度理解……",
        "提到{point}，关键在于……",
    ]

    # 数据解释模板
    DATA_TEMPLATES = [
        "从数据来看，{data}反映了……",
        "这个数据{data}说明了……",
        "值得注意的是{data}，它表明……",
    ]

    # 结束语模板
    CLOSING_TEMPLATES = [
        "以上就是今天的分享，感谢聆听！",
        "感谢各位的耐心聆听，谢谢！",
        "我的分享到此结束，谢谢大家！",
    ]

    @staticmethod
    def generate_opening(title: str) -> str:
        import random
        template = random.choice(SpeakerNotesGenerator.OPENING_TEMPLATES)
        return template.format(title=title)

    @staticmethod
    def generate_transition(section_title: str) -> str:
        import random
        template = random.choice(SpeakerNotesGenerator.TRANSITION_TEMPLATES)
        return template.format(section=section_title)

    @staticmethod
    def generate_expansion(point: str) -> str:
        import random
        template = random.choice(SpeakerNotesGenerator.EXPANSION_TEMPLATES)
        return template.format(point=point)

    @staticmethod
    def generate_data_interpretation(data: str) -> str:
        import random
        template = random.choice(SpeakerNotesGenerator.DATA_TEMPLATES)
        return template.format(data=data)

    @staticmethod
    def generate_closing() -> str:
        import random
        return random.choice(SpeakerNotesGenerator.CLOSING_TEMPLATES)

    @staticmethod
    def generate_notes_for_section(section: Section, prev_section: Section = None, next_section: Section = None) -> Dict:
        """为单个章节生成完整备注"""
        notes = {
            "opening": "",
            "transitions": [],
            "expansions": [],
            "data_interpretations": [],
            "closing": "",
        }

        # 开场白（仅第一个章节）
        if prev_section is None:
            notes["opening"] = SpeakerNotesGenerator.generate_opening(section.title)

        # 过渡语
        if prev_section:
            notes["transitions"].append(
                SpeakerNotesGenerator.generate_transition(section.title)
            )

        # 要点拓展（最多3条）
        for point in section.content[:3]:
            notes["expansions"].append(
                SpeakerNotesGenerator.generate_expansion(point)
            )

        # 数据解释
        if section.data:
            for key, value in section.data.items():
                notes["data_interpretations"].append(
                    SpeakerNotesGenerator.generate_data_interpretation(f"{key}: {value}")
                )

        # 结束语（仅最后一个章节）
        if next_section is None:
            notes["closing"] = SpeakerNotesGenerator.generate_closing()

        return notes

    @staticmethod
    def generate_all_notes(document: Document) -> List[Dict]:
        """为整个文档生成备注"""
        all_notes = []
        for i, section in enumerate(document.sections):
            prev_section = document.sections[i - 1] if i > 0 else None
            next_section = document.sections[i + 1] if i < len(document.sections) - 1 else None
            notes = SpeakerNotesGenerator.generate_notes_for_section(section, prev_section, next_section)
            all_notes.append({
                "section": section.title,
                "notes": notes,
            })
        return all_notes


# ==================== 动画/切换建议引擎 ====================

class AnimationSuggester:
    """动画/切换效果建议（规则引擎）"""

    # 内容类型 → 动画建议
    ANIMATION_RULES = [
        {
            "condition": lambda s: len(s.content) >= 3 and all(len(c) < 20 for c in s.content[:3]),
            "animation": "逐条出现",
            "transition": "淡入",
            "reason": "短要点列表适合逐条展示，避免观众提前读完",
        },
        {
            "condition": lambda s: any(kw in s.title for kw in ["数据", "统计", "趋势", "分析"]),
            "animation": "动态图表",
            "transition": "推入",
            "reason": "数据类内容适合动态图表展示，增强视觉冲击",
        },
        {
            "condition": lambda s: any(kw in s.title for kw in ["对比", "比较", "vs"]),
            "animation": "左右对比",
            "transition": "分割",
            "reason": "对比类内容适合左右分屏对比展示",
        },
        {
            "condition": lambda s: any(kw in s.title for kw in ["总结", "结论", "展望"]),
            "animation": "放大强调",
            "transition": "缩放",
            "reason": "总结类内容适合放大强调，引起观众注意",
        },
        {
            "condition": lambda s: any(kw in s.title for kw in ["流程", "步骤", "阶段"]),
            "animation": "顺序出现",
            "transition": "推进",
            "reason": "流程类内容适合按步骤顺序展示",
        },
        {
            "condition": lambda s: True,  # 默认
            "animation": "淡入",
            "transition": "淡出",
            "reason": "通用淡入效果，简洁大方",
        },
    ]

    @staticmethod
    def suggest(section: Section) -> Dict:
        """为单个章节建议动画效果"""
        for rule in AnimationSuggester.ANIMATION_RULES:
            if rule["condition"](section):
                return {
                    "animation": rule["animation"],
                    "transition": rule["transition"],
                    "reason": rule["reason"],
                }
        return {"animation": "无", "transition": "无", "reason": ""}

    @staticmethod
    def suggest_all(document: Document) -> List[Dict]:
        """为整个文档生成动画建议"""
        results = []
        for section in document.sections:
            suggestion = AnimationSuggester.suggest(section)
            results.append({
                "section": section.title,
                "animation": suggestion["animation"],
                "transition": suggestion["transition"],
                "reason": suggestion["reason"],
            })
        return results


# ==================== 配色方案引擎 ====================

class ColorSchemeEngine:
    """
    配色方案自动适配
    支持：企业品牌色输入 → 自动生成完整配色方案
    """

    @staticmethod
    def hex_to_rgb(hex_color: str) -> Tuple[int, int, int]:
        """HEX 颜色转 RGB"""
        hex_color = hex_color.lstrip("#")
        if len(hex_color) == 3:
            hex_color = "".join([c * 2 for c in hex_color])
        return tuple(int(hex_color[i:i + 2], 16) for i in (0, 2, 4))

    @staticmethod
    def rgb_to_hex(r: int, g: int, b: int) -> str:
        """RGB 转 HEX"""
        return f"#{r:02X}{g:02X}{b:02X}"

    @staticmethod
    def generate_complementary(hex_color: str) -> List[str]:
        """生成互补色方案"""
        r, g, b = ColorSchemeEngine.hex_to_rgb(hex_color)
        # 互补色（色轮对面）
        comp_r, comp_g, comp_b = 255 - r, 255 - g, 255 - b
        return [
            hex_color,
            ColorSchemeEngine.rgb_to_hex(comp_r, comp_g, comp_b),
            ColorSchemeEngine.rgb_to_hex(min(255, r + 50), min(255, g + 50), min(255, b + 50)),
            ColorSchemeEngine.rgb_to_hex(max(0, r - 30), max(0, g - 30), max(0, b - 30)),
        ]

    @staticmethod
    def generate_analogous(hex_color: str) -> List[str]:
        """生成类似色方案"""
        r, g, b = ColorSchemeEngine.hex_to_rgb(hex_color)
        colors = [hex_color]
        for offset in [30, 60, -30]:
            nr = min(255, max(0, r + offset))
            ng = min(255, max(0, g + offset // 2))
            nb = min(255, max(0, b + offset // 3))
            colors.append(ColorSchemeEngine.rgb_to_hex(nr, ng, nb))
        return colors

    @staticmethod
    def generate_triadic(hex_color: str) -> List[str]:
        """生成三角色方案"""
        r, g, b = ColorSchemeEngine.hex_to_rgb(hex_color)
        return [
            hex_color,
            ColorSchemeEngine.rgb_to_hex(g, b, r),
            ColorSchemeEngine.rgb_to_hex(b, r, g),
            ColorSchemeEngine.rgb_to_hex(max(0, r - 50), max(0, g - 50), max(0, b - 50)),
        ]

    @staticmethod
    def generate_scheme(brand_color: str, scheme_type: str = "complementary") -> Dict:
        """
        生成完整配色方案

        Args:
            brand_color: 企业品牌色（HEX 格式，如 #1A73E8）
            scheme_type: 配色类型（complementary/analogous/triadic）
        """
        if not brand_color.startswith("#"):
            brand_color = f"#{brand_color}"

        if scheme_type == "complementary":
            colors = ColorSchemeEngine.generate_complementary(brand_color)
        elif scheme_type == "analogous":
            colors = ColorSchemeEngine.generate_analogous(brand_color)
        elif scheme_type == "triadic":
            colors = ColorSchemeEngine.generate_triadic(brand_color)
        else:
            colors = ColorSchemeEngine.generate_complementary(brand_color)

        return {
            "primary": colors[0],       # 主色（品牌色）
            "secondary": colors[1],     # 辅助色
            "accent": colors[2],        # 强调色
            "background": colors[3],    # 背景色
            "scheme_type": scheme_type,
        }


# ==================== 图表推荐引擎 ====================

class ChartRecommender:
    """图表类型推荐（规则引擎）"""

    CHART_RULES = [
        {
            "condition": lambda s: any(kw in s.title for kw in ["趋势", "变化", "增长", "下降", "走势"]),
            "chart_type": "line",
            "reason": "趋势类数据最适合用折线图展示",
        },
        {
            "condition": lambda s: any(kw in s.title for kw in ["占比", "比例", "份额", "分布", "构成"]),
            "chart_type": "pie",
            "reason": "占比类数据最适合用饼图展示",
        },
        {
            "condition": lambda s: any(kw in s.title for kw in ["对比", "比较", "排名", "top"]),
            "chart_type": "bar",
            "reason": "对比类数据最适合用柱状图展示",
        },
        {
            "condition": lambda s: any(kw in s.title for kw in ["流程", "步骤", "阶段"]),
            "chart_type": "funnel",
            "reason": "流程类数据最适合用漏斗图展示",
        },
        {
            "condition": lambda s: True,  # 默认
            "chart_type": "bar",
            "reason": "默认使用柱状图，通用性最强",
        },
    ]

    @staticmethod
    def recommend(section: Section) -> Dict:
        """为单个章节推荐图表类型"""
        for rule in ChartRecommender.CHART_RULES:
            if rule["condition"](section):
                return {
                    "chart_type": rule["chart_type"],
                    "reason": rule["reason"],
                }
        return {"chart_type": "bar", "reason": "默认柱状图"}

    @staticmethod
    def recommend_all(document: Document) -> List[Dict]:
        """为整个文档生成图表推荐"""
        results = []
        for section in document.sections:
            rec = ChartRecommender.recommend(section)
            results.append({
                "section": section.title,
                "chart_type": rec["chart_type"],
                "reason": rec["reason"],
            })
        return results


# ==================== 排练辅助 ====================

class RehearsalHelper:
    """
    排练辅助：估算演讲时间，标注时间紧张页面
    语速假设：中文约 200 字/分钟
    """

    WORDS_PER_MINUTE = 200  # 中文语速

    @staticmethod
    def estimate_time(text: str) -> float:
        """估算演讲时间（分钟）"""
        # 统计中文字符数
        chinese_chars = len(re.findall(r"[\u4e00-\u9fff]", text))
        # 统计英文单词数
        english_words = len(re.findall(r"[a-zA-Z]+", text))
        total_words = chinese_chars + english_words
        return round(total_words / RehearsalHelper.WORDS_PER_MINUTE, 1)

    @staticmethod
    def analyze_pace(document: Document, notes: List[Dict] = None) -> Dict:
        """分析整个文档的演讲节奏"""
        page_times = []
        total_time = 0

        for i, section in enumerate(document.sections):
            # 计算该章节的演讲内容字数
            content_text = " ".join(section.content)
            if notes and i < len(notes):
                note_text = " ".join(notes[i].get("notes", {}).get("expansions", []))
                content_text += " " + note_text

            time_min = RehearsalHelper.estimate_time(content_text)
            page_times.append({
                "section": section.title,
                "time_minutes": time_min,
                "word_count": len(content_text),
            })
            total_time += time_min

        # 标注时间紧张的页面（超过 3 分钟）
        tight_pages = [p for p in page_times if p["time_minutes"] > 3]

        return {
            "total_time_minutes": round(total_time, 1),
            "page_count": len(page_times),
            "pages": page_times,
            "tight_pages": tight_pages,
            "suggestion": f"总时长约 {round(total_time, 1)} 分钟，其中 {len(tight_pages)} 个页面时间偏长，建议精简内容或拆分页面。",
        }


# ==================== PPT 渲染器 ====================

class PPTRenderer:
    """PPT 渲染器（python-pptx）"""

    # 内置配色方案
    THEMES = {
        "business": {
            "name": "商务蓝",
            "primary": (0x1A, 0x73, 0xE8),
            "secondary": (0x34, 0x98, 0xDB),
            "accent": (0xF3, 0x9C, 0x12),
            "background": (0xFF, 0xFF, 0xFF),
            "text": (0x33, 0x33, 0x33),
            "light": (0xEB, 0xF5, 0xFB),
        },
        "tech": {
            "name": "科技紫",
            "primary": (0x6C, 0x5C, 0xE7),
            "secondary": (0x0A, 0x85, 0xFF),
            "accent": (0x00, 0xD4, 0xAA),
            "background": (0x1A, 0x1A, 0x2E),
            "text": (0xFF, 0xFF, 0xFF),
            "light": (0x2D, 0x2D, 0x44),
        },
        "minimal": {
            "name": "极简灰",
            "primary": (0x2C, 0x3E, 0x50),
            "secondary": (0x7F, 0x8C, 0x8D),
            "accent": (0xE7, 0x4C, 0x3C),
            "background": (0xFA, 0xFA, 0xFA),
            "text": (0x2C, 0x3E, 0x50),
            "light": (0xEC, 0xF0, 0xF1),
        },
    }

    @staticmethod
    def render(document: Document, output_path: str, theme: str = "business",
               color_scheme: Dict = None, notes: List[Dict] = None,
               animations: List[Dict] = None) -> Dict:
        """
        渲染 PPT

        Args:
            document: 统一文档模型
            output_path: 输出文件路径
            theme: 主题名称（business/tech/minimal）
            color_scheme: 自定义配色方案（可选）
            notes: 演讲者备注（可选）
            animations: 动画建议（可选）
        """
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN
        except ImportError:
            return {"success": False, "error": "需要安装 python-pptx：pip install python-pptx"}

        # 获取配色
        if color_scheme:
            colors = {
                "primary": ColorSchemeEngine.hex_to_rgb(color_scheme.get("primary", "#1A73E8")),
                "secondary": ColorSchemeEngine.hex_to_rgb(color_scheme.get("secondary", "#3498DB")),
                "accent": ColorSchemeEngine.hex_to_rgb(color_scheme.get("accent", "#F39C12")),
                "background": ColorSchemeEngine.hex_to_rgb(color_scheme.get("background", "#FFFFFF")),
                "text": ColorSchemeEngine.hex_to_rgb(color_scheme.get("text", "#333333")),
                "light": ColorSchemeEngine.hex_to_rgb(color_scheme.get("light", "#EBF5FB")),
            }
        else:
            colors = PPTRenderer.THEMES.get(theme, PPTRenderer.THEMES["business"])

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # ===== 封面页 =====
        slide_layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(slide_layout)

        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*colors["background"])

        # 主色块
        shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(2.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*colors["primary"])
        shape.line.fill.background()

        # 标题
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = document.title or "演示文稿"
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.alignment = PP_ALIGN.CENTER

        # 副标题
        txBox2 = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11), Inches(1))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = f"共 {len(document.sections)} 个章节 | {sum(len(s.content) for s in document.sections)} 条要点"
        p2.font.size = Pt(16)
        p2.font.color.rgb = RGBColor(*colors["text"])
        p2.alignment = PP_ALIGN.CENTER

        # 日期
        txBox3 = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(11), Inches(0.5))
        tf3 = txBox3.text_frame
        p3 = tf3.paragraphs[0]
        p3.text = datetime.now().strftime("%Y年%m月%d日")
        p3.font.size = Pt(12)
        p3.font.color.rgb = RGBColor(*colors["secondary"])
        p3.alignment = PP_ALIGN.CENTER

        # ===== 目录页 =====
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*colors["background"])

        shape = slide.shapes.add_shape(1, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*colors["primary"])
        shape.line.fill.background()
        txBox = shape.text_frame
        p = txBox.paragraphs[0]
        p.text = "目录"
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

        toc_items = [f"{i+1}. {s.title}" for i, s in enumerate(document.sections)]
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11), Inches(5))
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, item in enumerate(toc_items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = item
            p.font.size = Pt(16)
            p.font.color.rgb = RGBColor(*colors["text"])
            p.space_after = Pt(8)

        # ===== 章节页 =====
        for idx, section in enumerate(document.sections):
            slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(slide_layout)
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(*colors["background"])

            # 章节标题背景
            shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.3), Inches(7.5))
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(*colors["primary"])
            shape.line.fill.background()

            # 章节编号
            shape = slide.shapes.add_shape(1, Inches(0.5), Inches(0.3), Inches(1.2), Inches(0.8))
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(*colors["primary"])
            shape.line.fill.background()
            txBox = shape.text_frame
            p = txBox.paragraphs[0]
            p.text = f"{idx+1}"
            p.font.size = Pt(24)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            p.alignment = PP_ALIGN.CENTER

            # 章节标题
            txBox = slide.shapes.add_textbox(Inches(2), Inches(0.3), Inches(10), Inches(0.8))
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = section.title
            p.font.size = Pt(24)
            p.font.bold = True
            p.font.color.rgb = RGBColor(*colors["primary"])

            # 章节内容
            content = section.content
            if content:
                bullets = [c for c in content if len(c) > 5][:8]
                txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11), Inches(5.5))
                tf = txBox.text_frame
                tf.word_wrap = True
                for i, bullet in enumerate(bullets):
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    text = bullet if len(bullet) <= 80 else bullet[:77] + "..."
                    p.text = f"• {text}"
                    p.font.size = Pt(14)
                    p.font.color.rgb = RGBColor(*colors["text"])
                    p.space_after = Pt(6)

            # 页码
            txBox = slide.shapes.add_textbox(Inches(12), Inches(7), Inches(1), Inches(0.3))
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = f"{idx+1} / {len(document.sections)}"
            p.font.size = Pt(10)
            p.font.color.rgb = RGBColor(*colors["secondary"])
            p.alignment = PP_ALIGN.RIGHT

        # ===== 结束页 =====
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*colors["primary"])

        txBox = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(9), Inches(1.5))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = "谢谢"
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.alignment = PP_ALIGN.CENTER

        # 保存
        Path(output_path).parent.mkdir(parents=True, exist_ok=True)
        prs.save(output_path)

        return {
            "success": True,
            "path": output_path,
            "title": document.title,
            "theme": theme,
            "total_sections": len(document.sections),
            "total_slides": len(document.sections) + 3,
            "engine": "PURE",
        }


# ==================== 通用入口 ====================

def generate_ppt(input_path: str, output_path: str = None,
                 input_type: str = "auto", title: str = None,
                 theme: str = "business", brand_color: str = None,
                 scheme_type: str = "complementary",
                 enable_notes: bool = True,
                 enable_animations: bool = True,
                 enable_rehearsal: bool = True) -> Dict:
    """
    通用 PPT 生成入口

    Args:
        input_path: 输入文件路径（Word/Markdown/思维导图/纯文本）
        output_path: 输出 PPT 路径（可选，默认同名 .pptx）
        input_type: 输入类型（auto/word/markdown/mindmap/text）
        title: PPT 标题（可选，默认取文件名）
        theme: 主题名称（business/tech/minimal）
        brand_color: 企业品牌色（HEX 格式，可选）
        scheme_type: 配色类型（complementary/analogous/triadic）
        enable_notes: 是否生成演讲者备注
        enable_animations: 是否生成动画建议
        enable_rehearsal: 是否生成排练报告
    """
    try:
        # 1. 解析输入
        document = parse_input(input_path, input_type)

        if not document.sections:
            return {"success": False, "error": "文档中未找到有效章节结构"}

        # 2. 设置标题
        if title:
            document.title = title
        elif not document.title:
            document.title = Path(input_path).stem

        # 3. 生成输出路径
        out_path = output_path or str(Path(input_path).with_suffix(".pptx"))
        if not out_path.endswith(".pptx"):
            out_path += ".pptx"

        # 4. 生成配色方案
        color_scheme = None
        if brand_color:
            color_scheme = ColorSchemeEngine.generate_scheme(brand_color, scheme_type)

        # 5. 生成演讲者备注
        notes = None
        if enable_notes:
            notes = SpeakerNotesGenerator.generate_all_notes(document)

        # 6. 生成动画建议
        animations = None
        if enable_animations:
            animations = AnimationSuggester.suggest_all(document)

        # 7. 渲染 PPT
        result = PPTRenderer.render(document, out_path, theme, color_scheme, notes, animations)

        if not result["success"]:
            return result

        # 8. 生成排练报告
        rehearsal = None
        if enable_rehearsal:
            rehearsal = RehearsalHelper.analyze_pace(document, notes)

        # 9. 组装结果
        result["document"] = document.to_dict()
        result["notes"] = notes
        result["animations"] = animations
        result["rehearsal"] = rehearsal
        result["color_scheme"] = color_scheme

        return result

    except Exception as e:
        return {"success": False, "error": f"生成失败: {str(e)}"}


# ==================== CLI ====================

if __name__ == "__main__":
    import argparse

    parser = argparse.ArgumentParser(description="PPT 智能生成器 v4.1")
    sub = parser.add_subparsers(dest="command", required=True)

    p_generate = sub.add_parser("generate", help="生成 PPT")
    p_generate.add_argument("--input", required=True, help="输入文件路径")
    p_generate.add_argument("--output", default="", help="输出 PPT 路径")
    p_generate.add_argument("--type", choices=["auto", "word", "markdown", "mindmap", "text"], default="auto")
    p_generate.add_argument("--title", default="", help="PPT 标题")
    p_generate.add_argument("--theme", choices=["business", "tech", "minimal"], default="business")
    p_generate.add_argument("--brand-color", default="", help="企业品牌色（HEX 格式，如 #1A73E8）")
    p_generate.add_argument("--scheme-type", choices=["complementary", "analogous", "triadic"], default="complementary")
    p_generate.add_argument("--no-notes", action="store_true", help="不生成演讲者备注")
    p_generate.add_argument("--no-animations", action="store_true", help="不生成动画建议")
    p_generate.add_argument("--no-rehearsal", action="store_true", help="不生成排练报告")

    p_scheme = sub.add_parser("scheme", help="生成配色方案")
    p_scheme.add_argument("--color", required=True, help="品牌色（HEX 格式）")
    p_scheme.add_argument("--type", choices=["complementary", "analogous", "triadic"], default="complementary")

    args = parser.parse_args()

    if args.command == "generate":
        result = generate_ppt(
            input_path=args.input,
            output_path=args.output,
            input_type=args.type,
            title=args.title,
            theme=args.theme,
            brand_color=args.brand_color,
            scheme_type=args.scheme_type,
            enable_notes=not args.no_notes,
            enable_animations=not args.no_animations,
            enable_rehearsal=not args.no_rehearsal,
        )
    elif args.command == "scheme":
        result = ColorSchemeEngine.generate_scheme(args.color, args.type)
    else:
        result = {"success": False, "error": "未知命令"}

    print(json.dumps(result, ensure_ascii=False, default=str))
