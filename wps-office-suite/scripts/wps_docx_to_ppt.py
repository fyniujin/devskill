"""
Word → PPT 一键生成器 v4.0
输入 .docx 文件或讲稿文字 → 自动分析章节结构 → 生成大纲 → 选模板 → 输出完整 PPT

核心能力：
  • 章节分割（基于 word 标题样式或文本特征）
  • 智能大纲生成（层级识别）
  • 3 种内置模板（商务/科技/极简）
  • 纯 Python 实现（python-pptx），无需 WPS

不依赖任何外部 API，完全本地处理。
"""
import re
import json
from pathlib import Path
from typing import List, Dict, Optional, Tuple


# ==================== 章节分割引擎 ====================

def split_chapters_by_style(filepath: str) -> List[Dict]:
    """基于 Word 标题样式分割章节"""
    try:
        from docx import Document
        doc = Document(filepath)
        chapters = []
        current_chapter = None

        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue

            style_name = para.style.name if para.style else ""
            is_heading = style_name.startswith("Heading") or style_name.startswith("标题")

            if is_heading:
                match = re.search(r"(\d+)", style_name)
                level = int(match.group(1)) if match else 1

                if current_chapter:
                    chapters.append(current_chapter)

                current_chapter = {
                    "title": text,
                    "level": level,
                    "content": [],
                    "sub_chapters": []
                }
            elif current_chapter is not None:
                current_chapter["content"].append(text)

        if current_chapter:
            chapters.append(current_chapter)

        return chapters
    except Exception as e:
        return []


def split_chapters_by_text(text: str) -> List[Dict]:
    """基于文本特征分割章节（无 Word 样式时使用）"""
    lines = text.strip().split("\n")
    chapters = []
    current_chapter = None

    for line in lines:
        line = line.strip()
        if not line:
            continue

        # 检测标题特征
        is_heading = False
        level = 1

        # 数字编号开头：1. 或 1、或 第一章
        if re.match(r"^第[一二三四五六七八九十百千]+[章节部分]", line):
            is_heading = True
            level = 1
        elif re.match(r"^\d+[\.\、]", line):
            is_heading = True
            level = 1
        elif re.match(r"^[\(（]\d+[\)）]", line):
            is_heading = True
            level = 2
        elif line.startswith("#") or line.startswith("【"):
            is_heading = True
            level = 1
        elif len(line) < 30 and not line.endswith(("，", "。", "？", "！")):
            # 短行无标点，可能是标题
            is_heading = True
            level = 2

        if is_heading:
            if current_chapter:
                chapters.append(current_chapter)
            current_chapter = {
                "title": line,
                "level": level,
                "content": [],
                "sub_chapters": []
            }
        elif current_chapter is not None:
            current_chapter["content"].append(line)

    if current_chapter:
        chapters.append(current_chapter)

    return chapters


def analyze_structure(chapters: List[Dict]) -> Dict:
    """分析文档结构，生成大纲"""
    outline = {
        "title": "演示大纲",
        "total_chapters": len(chapters),
        "levels": 0,
        "tree": []
    }

    max_level = 0
    for ch in chapters:
        if ch["level"] > max_level:
            max_level = ch["level"]

        # 计算内容字数
        word_count = sum(len(c) for c in ch["content"])

        node = {
            "title": ch["title"],
            "level": ch["level"],
            "word_count": word_count,
            "has_content": len(ch["content"]) > 0
        }
        outline["tree"].append(node)

    outline["levels"] = max_level
    return outline


# ==================== PPT 模板引擎 ====================

# 内置配色方案
THEMES = {
    "business": {
        "name": "商务蓝",
        "primary": (0x1A, 0x73, 0xE8),
        "secondary": (0x34, 0x98, 0xDB),
        "accent": (0xF3, 0x9C, 0x12),
        "background": (0xFF, 0xFF, 0xFF),
        "text": (0x33, 0x33, 0x33),
        "light": (0xEB, 0xF5, 0xFB)
    },
    "tech": {
        "name": "科技紫",
        "primary": (0x6C, 0x5C, 0xE7),
        "secondary": (0xA, 0x85, 0xFF),
        "accent": (0x00, 0xD4, 0xAA),
        "background": (0x1A, 0x1A, 0x2E),
        "text": (0xFF, 0xFF, 0xFF),
        "light": (0x2D, 0x2D, 0x44)
    },
    "minimal": {
        "name": "极简灰",
        "primary": (0x2C, 0x3E, 0x50),
        "secondary": (0x7F, 0x8C, 0x8D),
        "accent": (0xE7, 0x4C, 0x3C),
        "background": (0xFA, 0xFA, 0xFA),
        "text": (0x2C, 0x3E, 0x50),
        "light": (0xEC, 0xF0, 0xF1)
    }
}


def generate_ppt_from_chapters(chapters: List[Dict], output_path: str,
                                title: str = "演示文稿",
                                theme: str = "business") -> Dict:
    """
    将章节列表转换为 PPT
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        return {"success": False, "error": "需要安装 python-pptx：pip install python-pptx"}

    if theme not in THEMES:
        theme = "business"

    colors = THEMES[theme]
    prs = Presentation()

    # 设置幻灯片尺寸（16:9）
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ===== 封面页 =====
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # 背景色
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*colors["background"])

    # 主色块
    from pptx.util import Emu
    shape = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0), Inches(0),
        Inches(13.333), Inches(2.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*colors["primary"])
    shape.line.fill.background()

    # 标题文本
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.CENTER

    # 副标题
    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11), Inches(1))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = f"共 {len(chapters)} 个章节 | {sum(len(ch['content']) for ch in chapters)} 条要点"
    p2.font.size = Pt(16)
    p2.font.color.rgb = RGBColor(*colors["text"])
    p2.alignment = PP_ALIGN.CENTER

    # 日期
    txBox3 = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(11), Inches(0.5))
    tf3 = txBox3.text_frame
    from datetime import datetime
    p3 = tf3.paragraphs[0]
    p3.text = datetime.now().strftime("%Y年%m月%d日")
    p3.font.size = Pt(12)
    p3.font.color.rgb = RGBColor(*colors["secondary"])
    p3.alignment = PP_ALIGN.CENTER

    # ===== 目录页 =====
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*colors["background"])

    # 目录标题
    shape = slide.shapes.add_shape(1, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*colors["primary"])
    shape.line.fill.background()
    txBox = shape.text_frame
    p = txBox.paragraphs[0]
    p.text = "目录"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # 目录内容
    toc_items = [f"{i+1}. {ch['title']}" for i, ch in enumerate(chapters)]
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(toc_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(*colors["text"])
        p.space_after = Pt(8)

    # ===== 章节页 =====
    for idx, chapter in enumerate(chapters):
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*colors["background"])

        # 章节标题背景
        shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.3), Inches(7.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*colors["primary"])
        shape.line.fill.background()

        # 章节编号
        shape = slide.shapes.add_shape(1, Inches(0.5), Inches(0.3), Inches(1.2), Inches(0.8))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*colors["primary"])
        shape.line.fill.background()
        txBox = shape.text_frame
        p = txBox.paragraphs[0]
        p.text = f"{idx+1}"
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.alignment = PP_ALIGN.CENTER

        # 章节标题
        txBox = slide.shapes.add_textbox(Inches(2), Inches(0.3), Inches(10), Inches(0.8))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = chapter["title"]
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = RGBColor(*colors["primary"])

        # 章节内容
        content = chapter.get("content", [])
        if content:
            # 只取前8条要点（避免页面拥挤）
            bullets = [c for c in content if len(c) > 5][:8]
            txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11), Inches(5.5))
            tf = txBox.text_frame
            tf.word_wrap = True

            for i, bullet in enumerate(bullets):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                # 截断过长内容
                text = bullet if len(bullet) <= 80 else bullet[:77] + "..."
                p.text = f"• {text}"
                p.font.size = Pt(14)
                p.font.color.rgb = RGBColor(*colors["text"])
                p.space_after = Pt(6)

        # 页码
        txBox = slide.shapes.add_textbox(Inches(12), Inches(7), Inches(1), Inches(0.3))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = f"{idx+1} / {len(chapters)}"
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(*colors["secondary"])
        p.alignment = PP_ALIGN.RIGHT

    # ===== 结束页 =====
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*colors["primary"])

    txBox = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(9), Inches(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "谢谢"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.CENTER

    # 保存
    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)

    return {
        "success": True,
        "path": output_path,
        "title": title,
        "theme": theme,
        "total_chapters": len(chapters),
        "total_slides": len(chapters) + 3,  # 封面 + 目录 + 章节 + 结束
        "engine": "PURE"
    }


# ==================== 通用入口 ====================

def docx_to_ppt(input_path: str, output_path: str = None,
                title: str = None, theme: str = "business") -> Dict:
    """
    将 Word 文档转换为 PPT

    Args:
        input_path: Word 文档路径（.docx）
        output_path: 输出 PPT 路径（可选，默认同名 .pptx）
        title: PPT 标题（可选，默认取文件名）
        theme: 模板名称（business/tech/minimal）

    Returns:
        dict: 转换结果
    """
    try:
        input_p = Path(input_path).resolve()
        if not input_p.exists():
            return {"success": False, "error": f"文件不存在: {input_path}"}

        if not input_p.suffix.lower() == ".docx":
            return {"success": False, "error": "输入必须是 .docx 文件"}

        # 分割章节
        chapters = split_chapters_by_style(str(input_p))

        if not chapters:
            return {"success": False, "error": "文档中未找到有效章节（Heading 样式）"}

        # 生成标题
        ppt_title = title or input_p.stem

        # 生成输出路径
        out_path = output_path or str(input_p.with_suffix(".pptx"))
        if not out_path.endswith(".pptx"):
            out_path += ".pptx"

        return generate_ppt_from_chapters(chapters, out_path, ppt_title, theme)

    except Exception as e:
        return {"success": False, "error": f"转换失败: {str(e)}"}


def text_to_ppt(text: str, output_path: str,
                title: str = "演示文稿", theme: str = "business") -> Dict:
    """
    将纯文本转换为 PPT

    Args:
        text: 讲稿文字内容
        output_path: 输出 PPT 路径
        title: PPT 标题
        theme: 模板名称
    """
    try:
        chapters = split_chapters_by_text(text)

        if not chapters:
            return {"success": False, "error": "文本中未找到有效章节结构"}

        return generate_ppt_from_chapters(chapters, output_path, title, theme)

    except Exception as e:
        return {"success": False, "error": f"转换失败: {str(e)}"}


# ==================== CLI ====================

if __name__ == "__main__":
    import argparse

    parser = argparse.ArgumentParser(description="Word → PPT 一键生成器")
    sub = parser.add_subparsers(dest="command", required=True)

    p_file = sub.add_parser("file", help="从 Word 文件生成")
    p_file.add_argument("--input", required=True, help="输入 .docx 文件路径")
    p_file.add_argument("--output", default="", help="输出 .pptx 文件路径")
    p_file.add_argument("--title", default="", help="PPT 标题")
    p_file.add_argument("--theme", choices=["business", "tech", "minimal"], default="business")

    p_text = sub.add_parser("text", help="从文本生成")
    p_text.add_argument("--text", required=True, help="讲稿文字内容")
    p_text.add_argument("--output", required=True, help="输出 .pptx 文件路径")
    p_text.add_argument("--title", default="演示文稿", help="PPT 标题")
    p_text.add_argument("--theme", choices=["business", "tech", "minimal"], default="business")

    args = parser.parse_args()

    if args.command == "file":
        result = docx_to_ppt(args.input, args.output, args.title, args.theme)
    else:
        result = text_to_ppt(args.text, args.output, args.title, args.theme)

    print(json.dumps(result, ensure_ascii=False, default=str))
