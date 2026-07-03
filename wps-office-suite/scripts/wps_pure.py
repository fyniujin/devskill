"""
纯 Python 模式 v2.5 - 零依赖办公文档操作引擎
支持功能：
  • 创建、编辑、格式设置 — Word/Excel/PPT
  • 表格创建、图片插入、列宽调整 — Word
  • 多列排序（升/降序）、多条件筛选（AND/OR）— Excel
  • 图表生成（柱/折/饼）— Excel
  • 数据汇总统计（SUM/AVG/COUNT/MAX/MIN）— Excel
  • 目录提取/生成 — Word
  • 回退 LibreOffice Headless — 格式转换
"""
from pathlib import Path
from typing import Any, Dict, List, Optional


def _cell_value(cell_val):
    """提取单元格原始值"""
    if cell_val is None:
        return None
    if isinstance(cell_val, (int, float)):
        return cell_val
    s = str(cell_val).strip()
    if s.startswith("=") or s.startswith("+"):
        return s
    try:
        return float(s) if "." in s else int(s)
    except (ValueError, TypeError):
        return s


def _compare(op: str, cell_val, target_val) -> bool:
    """
    比较操作。
    op: ">", "<", ">=", "<=", "=", "!=", "contains", "startswith", "endswith"
    """
    if op == "contains":
        return target_val is not None and str(target_val).lower() in str(cell_val).lower() if cell_val is not None else False
    if op == "startswith":
        return cell_val is not None and str(cell_val).lower().startswith(str(target_val).lower())
    if op == "endswith":
        return cell_val is not None and str(cell_val).lower().endswith(str(target_val).lower())
    # 数值比较
    if cell_val is None:
        if op == "=":
            return target_val is None
        if op == "!=":
            return target_val is not None
        return False
    try:
        c = float(cell_val)
        t = float(target_val)
        if op == ">":
            return c > t
        if op == "<":
            return c < t
        if op == ">=":
            return c >= t
        if op == "<=":
            return c <= t
        if op == "=":
            return c == t
        if op == "!=":
            return c != t
    except (ValueError, TypeError):
        c = str(cell_val)
        t = str(target_val)
        if op == "=":
            return c == t
        if op == "!=":
            return c != t
    return False


# ==================== Word 模块 ====================

def pure_create_word(title, filepath, body="",
                     font_name="微软雅黑", font_size=12,
                     align="left"):
    """创建 Word 文档"""
    try:
        from docx import Document
        from docx.shared import Pt, Cm, RGBColor, Inches
        from docx.enum.text import WD_ALIGN_PARAGRAPH

        doc = Document()
        style = doc.styles["Normal"]
        style.font.name = font_name
        style.font.size = Pt(font_size)

        h = doc.add_heading(title, level=1)
        h.alignment = WD_ALIGN_PARAGRAPH.CENTER
        for run in h.runs:
            run.font.color.rgb = RGBColor(0x1A, 0x73, 0xE8)

        if body:
            for line in body.split("\n"):
                if line.strip():
                    p = doc.add_paragraph(line.strip())
                    am = {"center": 1, "left": 0, "right": 2, "justify": 3}
                    if align in am:
                        p.alignment = WD_ALIGN_PARAGRAPH(am[align])
                    for run in p.runs:
                        run.font.name = font_name
                        run.font.size = Pt(font_size)

        Path(filepath).parent.mkdir(parents=True, exist_ok=True)
        doc.save(filepath)
        return {"success": True, "path": filepath, "engine": "PURE"}
    except ImportError:
        return {"success": False, "error": "需要安装 python-docx：pip install python-docx"}
    except Exception as e:
        return {"success": False, "error": str(e)}


def pure_edit_word(filepath, text, position="end"):
    """追加/前方插入文本"""
    try:
        from docx import Document
        if not Path(filepath).exists():
            return {"success": False, "error": f"文件不存在: {filepath}"}
        doc = Document(filepath)
        if position == "end":
            for line in text.split("\n"):
                if line.strip():
                    doc.add_paragraph(line.strip())
        elif position == "start":
            for line in reversed(text.split("\n")):
                if line.strip():
                    if doc.paragraphs:
                        doc.paragraphs[0].insert_paragraph_before(line.strip())
                    else:
                        doc.add_paragraph(line.strip())
        doc.save(filepath)
        return {"success": True, "engine": "PURE"}
    except Exception as e:
        return {"success": False, "error": str(e)}


def pure_format_word(filepath, font_name="", font_size=0, bold=False,
                    align="", space_after=0, first_line_indent=0,
                    color=""):
    """格式设置"""
    try:
        from docx import Document
        from docx.shared import Pt, Cm, RGBColor
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        if not Path(filepath).exists():
            return {"success": False, "error": f"文件不存在: {filepath}"}
        doc = Document(filepath)
        am = {"center": 1, "left": 0, "right": 2, "justify": 3}
        cm = {"red": (0xE7, 0x4C, 0x3C), "blue": (0x34, 0x98, 0xDB),
              "green": (0x27, 0xAE, 0x60), "black": (0x00, 0x00, 0x00),
              "gray": (0x95, 0xA5, 0xA6), "orange": (0xF3, 0x9C, 0x12)}
        rgb = cm.get(color.lower(), None)

        for para in doc.paragraphs:
            if not para.text.strip():
                continue
            if align in am:
                para.alignment = WD_ALIGN_PARAGRAPH(am[align])
            if space_after > 0:
                para.paragraph_format.space_after = Pt(space_after)
            if first_line_indent > 0:
                para.paragraph_format.first_line_indent = Cm(first_line_indent)
            for run in para.runs:
                if font_name:
                    run.font.name = font_name
                if font_size > 0:
                    run.font.size = Pt(font_size)
                if bold:
                    run.font.bold = True
                if rgb:
                    run.font.color.rgb = RGBColor(*rgb)

        doc.save(filepath)
        return {"success": True, "engine": "PURE", "formatted": True}
    except ImportError:
        return {"success": False, "error": "需要安装 python-docx"}
    except Exception as e:
        return {"success": False, "error": str(e)}


def pure_add_table(filepath, headers, rows):
    """在 Word 文档末尾添加表格"""
    try:
        from docx import Document
        from docx.shared import Pt, Cm
        from docx.enum.table import WD_TABLE_ALIGNMENT
        if not Path(filepath).exists():
            return {"success": False, "error": f"文件不存在: {filepath}"}
        doc = Document(filepath)
        table = doc.add_table(rows=1 + len(rows), cols=len(headers))
        table.style = "Table Grid"
        table.alignment = WD_TABLE_ALIGNMENT.CENTER

        for i, h in enumerate(headers):
            cell = table.rows[0].cells[i]
            cell.text = str(h)
            for para in cell.paragraphs:
                for run in para.runs:
                    run.font.bold = True
                    run.font.size = Pt(10)

        for r_idx, row in enumerate(rows):
            for c_idx, val in enumerate(row):
                if c_idx < len(headers):
                    table.rows[r_idx + 1].cells[c_idx].text = str(val)

        for col in table.columns:
            for cell in col.cells:
                cell.width = Cm(3)

        doc.save(filepath)
        return {"success": True, "engine": "PURE",
                "table_added": True, "rows": len(rows)}
    except ImportError:
        return {"success": False, "error": "需要安装 python-docx"}
    except Exception as e:
        return {"success": False, "error": str(e)}


def pure_insert_image(filepath, image_path, width_cm=10, caption=""):
    """插入图片"""
    try:
        from docx import Document
        from docx.shared import Cm
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        if not Path(filepath).exists():
            return {"success": False, "error": f"文件不存在: {filepath}"}
        if not Path(image_path).exists():
            return {"success": False, "error": f"图片不存在: {image_path}"}
        doc = Document(filepath)
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = p.add_run()
        run.add_picture(image_path, width=Cm(width_cm))
        if caption:
            cap = doc.add_paragraph(caption)
            cap.alignment = WD_ALIGN_PARAGRAPH.CENTER
            for run in cap.runs:
                run.font.size = Pt(9)
        doc.save(filepath)
        return {"success": True, "engine": "PURE", "image_inserted": True}
    except ImportError:
        return {"success": False, "error": "需要安装 python-docx"}
    except Exception as e:
        return {"success": False, "error": str(e)}


def pure_info_word(filepath):
    """Word 文档信息"""
    try:
        from docx import Document
        if not Path(filepath).exists():
            return {"success": False, "error": f"文件不存在: {filepath}"}
        doc = Document(filepath)
        stat = Path(filepath).stat()
        return {
            "success": True,
            "name": Path(filepath).name,
            "path": filepath,
            "size_bytes": stat.st_size,
            "paragraph_count": len(doc.paragraphs),
            "table_count": len(doc.tables),
            "engine": "PURE",
        }
    except ImportError:
        return {"success": False, "error": "需要安装 python-docx"}
    except Exception as e:
        return {"success": False, "error": str(e)}


# ==================== Excel 模块 ====================

def pure_create_excel(name, sheets, filepath, data=None):
    """创建 Excel"""
    try:
        import openpyxl
        from openpyxl.styles import Font, PatternFill, Alignment

        wb = openpyxl.Workbook()
        wb.remove(wb.active)

        for sheet_name in sheets:
            ws = wb.create_sheet(title=sheet_name)
            if data and sheet_name in data:
                sheet_data = data[sheet_name]
                for r, row in enumerate(sheet_data, start=1):
                    for c, val in enumerate(row, start=1):
                        cell = ws.cell(row=r, column=c, value=val)
                        if r == 1:
                            cell.font = Font(bold=True, color="FFFFFF")
                            cell.fill = PatternFill("solid", fgColor="1A73E8")
                            cell.alignment = Alignment(horizontal="center")

        Path(filepath).parent.mkdir(parents=True, exist_ok=True)
        wb.save(filepath)
        return {"success": True, "path": filepath, "sheets": sheets, "engine": "PURE"}
    except ImportError:
        return {"success": False, "error": "需要安装 openpyxl"}
    except Exception as e:
        return {"success": False, "error": str(e)}


def pure_input_excel(filepath, sheet, data):
    """录入数据"""
    try:
        import openpyxl
        if not Path(filepath).exists():
            return {"success": False, "error": f"文件不存在: {filepath}"}
        wb = openpyxl.load_workbook(filepath)
        if sheet not in wb.sheetnames:
            ws = wb.create_sheet(title=sheet)
        else:
            ws = wb[sheet]
        start_row = ws.max_row + 1
        for r_offset, row in enumerate(data):
            for c, val in enumerate(row, start=1):
                ws.cell(row=start_row + r_offset, column=c, value=val)
        wb.save(filepath)
        return {"success": True, "input_rows": len(data), "engine": "PURE"}
    except ImportError:
        return {"success": False, "error": "需要安装 openpyxl"}
    except Exception as e:
        return {"success": False, "error": str(e)}


def pure_add_excel_sheet(filepath, sheet_name, headers=None, data=None):
    """新增 Sheet（含表头和数据）"""
    try:
        import openpyxl
        if not Path(filepath).exists():
            return {"success": False, "error": f"文件不存在: {filepath}"}
        wb = openpyxl.load_workbook(filepath)
        if sheet_name in wb.sheetnames:
            del wb[sheet_name]
        ws = wb.create_sheet(title=sheet_name)
        if headers:
            for c, h in enumerate(headers, start=1):
                ws.cell(row=1, column=c, value=h)
        if data:
            for r, row in enumerate(data, start=2):
                for c, val in enumerate(row, start=1):
                    ws.cell(row=r, column=c, value=val)
        wb.save(filepath)
        return {"success": True, "sheet_added": sheet_name, "engine": "PURE"}
    except Exception as e:
        return {"success": False, "error": str(e)}


def pure_sort_excel(filepath, sheet, sorts):
    """
    纯Python多列排序。
    sorts: [{"column": "A"/"列名", "ascending": true/false}, ...]
    第一行作为表头不参与排序。
    """
    try:
        import openpyxl
        if not Path(filepath).exists():
            return {"success": False, "error": f"文件不存在: {filepath}"}

        wb = openpyxl.load_workbook(filepath, data_only=True)
        if sheet not in wb.sheetnames:
            return {"success": False, "error": f"Sheet不存在: {sheet}"}
        ws = wb[sheet]

        # 读取 header
        max_col = ws.max_column
        max_row = ws.max_row
        if max_row <= 1:
            return {"success": True, "note": "数据不足，无需排序", "engine": "PURE"}

        header = []
        for c in range(1, max_col + 1):
            val = ws.cell(row=1, column=c).value
            header.append(str(val).strip() if val is not None else "")

        # 读取所有 数据行
        body_rows = []
        for r in range(2, max_row + 1):
            row = []
            for c in range(1, max_col + 1):
                cell_val = ws.cell(row=r, column=c).value
                if isinstance(cell_val, str):
                    row.append(cell_val.strip())
                else:
                    row.append(cell_val)
            body_rows.append(row)

        # 列索引查找
        def col_to_idx(col_ref):
            col_ref_str = str(col_ref).strip()
            # 尝试作为列名
            for i, h in enumerate(header):
                if h.lower() == col_ref_str.lower():
                    return i
            # 尝试作为列号 (0-indexed or A=1)
            try:
                return int(col_ref_str) - 1
            except ValueError:
                pass
            # 尝试作为 Excel column letter A, B, C...
            col_ref_str = col_ref_str.upper()
            if len(col_ref_str) == 1 and "A" <= col_ref_str <= "Z":
                return ord(col_ref_str) - ord("A")
            return None

        # 构建排序键
        sort_keys = []
        for s in sorts:
            col_ref = s.get("column", "")
            ascending = s.get("ascending", True)
            idx = col_to_idx(col_ref)
            if idx is None:
                return {"success": False, "error": f"列'{col_ref}'未找到"}
            sort_keys.append((idx, ascending))

        def sort_key_func(row):
            key_parts = []
            for idx, ascending in sort_keys:
                val = row[idx] if idx < len(row) else None
                if val is None:
                    key_parts.append((2, ""))
                elif isinstance(val, (int, float)):
                    key_parts.append((0, val if ascending else -val))
                else:
                    try:
                        num = float(val)
                        key_parts.append((0, num if ascending else -num))
                    except (ValueError, TypeError):
                        key_parts.append((1, str(val).lower()))
            return tuple(key_parts)

        body_rows.sort(key=sort_key_func)

        # 重新写回
        wb2 = openpyxl.load_workbook(filepath)
        ws2 = wb2[sheet]
        for r, row in enumerate(body_rows, start=2):
            for c, val in enumerate(row, start=1):
                ws2.cell(row=r, column=c, value=val)

        wb2.save(filepath)
        return {"success": True, "sorts": sorts, "rows": len(body_rows), "engine": "PURE"}
    except ImportError:
        return {"success": False, "error": "需要安装 openpyxl"}
    except Exception as e:
        return {"success": False, "error": str(e)}


def pure_filter_excel(filepath, sheet, conditions, logic="AND"):
    """
    纯Python多条件筛选。
    conditions: [{"column": "A"/"列名", "op": ">", "value": "100"}, ...]
    logic: "AND" / "OR"
    """
    try:
        import openpyxl
        if not Path(filepath).exists():
            return {"success": False, "error": f"文件不存在: {filepath}"}

        wb = openpyxl.load_workbook(filepath, data_only=True)
        if sheet not in wb.sheetnames:
            return {"success": False, "error": f"Sheet不存在: {sheet}"}
        ws = wb[sheet]

        max_col = ws.max_column
        max_row = ws.max_row
        if max_row <= 1:
            return {"success": True, "filtered": 0, "rows": [], "engine": "PURE"}

        header = []
        for c in range(1, max_col + 1):
            val = ws.cell(row=1, column=c).value
            header.append(str(val).strip() if val is not None else "")

        body_rows = []
        for r in range(2, max_row + 1):
            row = []
            for c in range(1, max_col + 1):
                cell_val = ws.cell(row=r, column=c).value
                row.append(cell_val)
            body_rows.append(row)

        def col_to_idx(col_ref):
            col_ref_str = str(col_ref).strip()
            for i, h in enumerate(header):
                if h.lower() == col_ref_str.lower():
                    return i
            try:
                return int(col_ref_str) - 1
            except ValueError:
                pass
            col_ref_str = col_ref_str.upper()
            if len(col_ref_str) == 1 and "A" <= col_ref_str <= "Z":
                return ord(col_ref_str) - ord("A")
            return None

        # 解析条件
        parsed_conditions = []
        for cond in conditions:
            col_ref = cond.get("column", "")
            op = cond.get("op", "=")
            target = cond.get("value", "")
            idx = col_to_idx(col_ref)
            if idx is None:
                return {"success": False, "error": f"列'{col_ref}'未找到"}
            parsed_conditions.append((idx, op, target))

        # 筛选
        filtered = []
        for row in body_rows:
            if logic.upper() == "AND":
                match = all(_compare(op, row[idx] if idx < len(row) else None, target)
                            for idx, op, target in parsed_conditions)
            else:  # OR
                match = any(_compare(op, row[idx] if idx < len(row) else None, target)
                            for idx, op, target in parsed_conditions)
            if match:
                filtered.append(row)

        for c_idx, cell_val in enumerate(header):
            filtered.insert(c_idx, cell_val)

        return {"success": True, "filtered": len(filtered), "rows": filtered, "logic": logic, "engine": "PURE"}
    except ImportError:
        return {"success": False, "error": "需要安装 openpyxl"}
    except Exception as e:
        return {"success": False, "error": str(e)}


def pure_chart_excel(filepath, sheet, chart_type="bar", data_range="A1:B10",
                     title="", output_path=""):
    """
    纯 Python 图表生成。
    使用 openpyxl 图表模块直接生成嵌入式图表。
    支持: bar(柱), line(折), pie(饼)
    """
    try:
        import openpyxl
        from openpyxl.chart import BarChart, LineChart, PieChart, Reference

        if not Path(filepath).exists():
            return {"success": False, "error": f"文件不存在: {filepath}"}

        wb = openpyxl.load_workbook(filepath)
        if sheet not in wb.sheetnames:
            return {"success": False, "error": f"Sheet不存在: {sheet}"}
        ws = wb[sheet]

        # 解析 data_range
        parts = data_range.split(":")
        if len(parts) != 2:
            return {"success": False, "error": f"数据范围格式错误: {data_range}"}

        # 获取行列范围
        def parse_cell(cell_str):
            cell_str = cell_str.strip().upper()
            col_str = ""
            row_str = ""
            for ch in cell_str:
                if ch.isalpha():
                    col_str += ch
                else:
                    row_str += ch
            col = 0
            for c in col_str:
                col = col * 26 + (ord(c) - ord("A") + 1)
            return col, int(row_str) if row_str else 1

        col1, row1 = parse_cell(parts[0])
        col2, row2 = parse_cell(parts[1])

        if chart_type == "bar":
            chart = BarChart()
            chart.type = "col"
        elif chart_type == "line":
            chart = LineChart()
        elif chart_type == "pie":
            chart = PieChart()
        else:
            chart = BarChart()

        chart.title = title or f"图表-{chart_type}"
        chart.width = 15
        chart.height = 10

        # 数据引用
        data_ref = Reference(ws, min_col=col1, min_row=row1, max_col=col2, max_row=row2)
        # 类别（第一列）
        cats_ref = Reference(ws, min_col=col1, min_row=row1 + 1, max_row=row2)

        chart.add_data(data_ref, titles_from_data=True)
        chart.set_categories(cats_ref)
        ws.add_chart(chart, "E5")

        if output_path:
            wb.save(output_path)
        wb.save(filepath)
        return {"success": True, "chart_created": True, "engine": "PURE",
                "chart_type": chart_type, "data_range": data_range}
    except ImportError:
        return {"success": False, "error": "需要安装 openpyxl"}
    except Exception as e:
        return {"success": False, "error": str(e)}


def pure_formula_excel(filepath, sheet, cell, formula):
    """
    纯 Python 模式写入 Excel 公式。
    openpyxl 支持写入公式字符串，打开时 Excel/WPS 会自动计算。
    """
    try:
        import openpyxl
        if not Path(filepath).exists():
            return {"success": False, "error": f"文件不存在: {filepath}"}

        wb = openpyxl.load_workbook(filepath)
        if sheet not in wb.sheetnames:
            ws = wb.create_sheet(title=sheet)
        else:
            ws = wb[sheet]

        ws[cell] = formula
        wb.save(filepath)
        return {"success": True, "formula_set": True, "cell": cell, "formula": formula, "engine": "PURE"}
    except ImportError:
        return {"success": False, "error": "需要安装 openpyxl"}
    except Exception as e:
        return {"success": False, "error": str(e)}


def pure_info_excel(filepath):
    """Excel 文档信息"""
    try:
        import openpyxl
        if not Path(filepath).exists():
            return {"success": False, "error": f"文件不存在: {filepath}"}
        wb = openpyxl.load_workbook(filepath, read_only=True)
        stat = Path(filepath).stat()
        return {
            "success": True,
            "name": Path(filepath).name,
            "path": filepath,
            "sheets": wb.sheetnames,
            "size_bytes": stat.st_size,
            "sheet_count": len(wb.sheetnames),
            "engine": "PURE",
        }
    except ImportError:
        return {"success": False, "error": "需要安装 openpyxl"}
    except Exception as e:
        return {"success": False, "error": str(e)}


def pure_statistics_excel(filepath, sheet, column, stat_type="SUM"):
    """
    数据汇总统计（跨平台替代透视表）。
    stat_type: SUM, AVG, COUNT, MAX, MIN, COUNTA
    """
    try:
        import openpyxl
        if not Path(filepath).exists():
            return {"success": False, "error": f"文件不存在: {filepath}"}
        wb = openpyxl.load_workbook(filepath, data_only=True)
        if sheet not in wb.sheetnames:
            return {"success": False, "error": f"Sheet不存在: {sheet}"}
        ws = wb[sheet]

        max_col = ws.max_column
        header = []
        for c in range(1, max_col + 1):
            val = ws.cell(row=1, column=c).value
            header.append(str(val).strip() if val is not None else "")

        col_ref = str(column).strip()
        col_idx = None
        for i, h in enumerate(header):
            if h.lower() == col_ref.lower():
                col_idx = i + 1
                break
        if col_idx is None:
            try:
                col_ref_upper = col_ref.upper()
                if len(col_ref_upper) == 1 and "A" <= col_ref_upper <= "Z":
                    col_idx = ord(col_ref_upper) - ord("A") + 1
            except:
                pass
        if col_idx is None:
            return {"success": False, "error": f"列'{col_ref}'未找到"}

        values = []
        for r in range(2, ws.max_row + 1):
            val = ws.cell(row=r, column=col_idx).value
            if val is None:
                continue
            if stat_type in ("SUM", "AVG", "MAX", "MIN"):
                try:
                    values.append(float(val))
                except (ValueError, TypeError):
                    pass
            elif stat_type in ("COUNT", "COUNTA"):
                values.append(val)

        if stat_type == "SUM":
            result = sum(values) if values else 0
        elif stat_type == "AVG":
            result = sum(values) / len(values) if values else 0
        elif stat_type == "MAX":
            result = max(values) if values else 0
        elif stat_type == "MIN":
            result = min(values) if values else 0
        elif stat_type == "COUNT":
            result = len(values)
        elif stat_type == "COUNTA":
            result = len(values)
        else:
            return {"success": False, "error": f"统计类型不支持: {stat_type}"}

        return {"success": True, "stat_type": stat_type, "column": col_ref,
                "result": result, "count": len(values), "engine": "PURE"}

    except ImportError:
        return {"success": False, "error": "需要安装 openpyxl"}
    except Exception as e:
        return {"success": False, "error": str(e)}


# ==================== PPT 模块 ====================

def pure_create_ppt(title, filepath, slides_content=None):
    """创建 PPT（支持多页）"""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN

        prs = Presentation()
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        if slide.shapes.title:
            slide.shapes.title.text = title
            slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(32)
            slide.shapes.title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        if slides_content:
            for content in slides_content:
                slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(slide_layout)
                if slide.shapes.title:
                    slide.shapes.title.text = content.get("title", "")
                for placeholder in slide.placeholders:
                    if placeholder.placeholder_format.idx == 1:
                        placeholder.text = content.get("body", "")

        Path(filepath).parent.mkdir(parents=True, exist_ok=True)
        prs.save(filepath)
        return {"success": True, "path": filepath, "engine": "PURE",
                "slides": 1 + (len(slides_content) if slides_content else 0)}
    except ImportError:
        return {
            "success": False,
            "error": "需要安装 python-pptx：pip install python-pptx",
        }
    except Exception as e:
        return {"success": False, "error": str(e)}


def pure_info_ppt(filepath):
    """PPT 信息"""
    try:
        from pptx import Presentation
        if not Path(filepath).exists():
            return {"success": False, "error": f"文件不存在: {filepath}"}
        prs = Presentation(filepath)
        stat = Path(filepath).stat()
        return {
            "success": True,
            "name": Path(filepath).name,
            "path": filepath,
            "size_bytes": stat.st_size,
            "slide_count": len(prs.slides),
            "engine": "PURE",
        }
    except ImportError:
        return {"success": False, "error": "需要安装 python-pptx"}
    except Exception as e:
        return {"success": False, "error": str(e)}


# ==================== 安全路径 ====================

def pure_safe_path(path_str):
    """安全路径校验"""
    path = Path(path_str).resolve()
    parent = path.parent
    if not parent.exists():
        raise FileNotFoundError(f"目录不存在：{parent}")
    problematic = '<>"|?*' if __import__("platform").system() == "Windows" else ""
    for ch in problematic:
        if ch in path.name:
            raise ValueError(f"文件名包含非法字符 '{ch}': {path.name}")
    return path


# ==================== 通用入口 ====================

def get_engine_info():
    return {"engine": "PURE", "wps_available": False,
            "msoffice_available": False, "pure_available": True}
