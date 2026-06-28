"""
WPS Office 全家桶 - 纯 Python 模式（无需安装 WPS/MS Office）
使用 python-docx / openpyxl / python-pptx 实现基础文档操作

可用能力：
- Word: 创建文档、添加段落/标题、格式设置、保存
- Excel: 创建工作簿、数据录入、保存
- PPT: 创建演示文稿、添加幻灯片、保存
- 格式转换: 仅支持通过 WPS/MS 模式，纯 Python 模式仅创建

自动检测模式优先级：WPS > MS Office > Pure Python
"""
from pathlib import Path
from typing import Optional, List, Dict, Any


def _detect_engine() -> str:
    """自动检测可用的办公引擎"""
    # 1. 尝试 WPS
    try:
        import win32com.client
        wps = win32com.client.Dispatch("WPS.Application")
        _ = wps.Name  # 测试是否真的可用
        return "WPS"
    except Exception:
        pass

    # 2. 尝试 MS Office
    try:
        import win32com.client
        word = win32com.client.Dispatch("Word.Application")
        _ = word.Name
        return "MSOFFICE"
    except Exception:
        pass

    # 3. 纯 Python 模式
    return "PURE"


ENGINE = _detect_engine()
ENGINE_PURE = "PURE"


# ==================== Word 纯 Python 模式 ====================

def pure_create_word(title: str, filepath: str, body: str = "") -> dict:
    """用 python-docx 创建 Word 文档"""
    try:
        from docx import Document
        from docx.shared import Pt, Inches
        from docx.enum.text import WD_ALIGN_PARAGRAPH

        doc = Document()

        # 添加标题
        heading = doc.add_heading(title, level=1)
        heading.alignment = WD_ALIGN_PARAGRAPH.CENTER

        # 添加正文
        if body:
            for line in body.split("\n"):
                if line.strip():
                    doc.add_paragraph(line.strip())

        # 确保父目录存在
        Path(filepath).parent.mkdir(parents=True, exist_ok=True)
        doc.save(filepath)
        return {"success": True, "path": filepath, "engine": "PURE"}
    except Exception as e:
        return {"success": False, "error": str(e)}


def pure_edit_word(filepath: str, text: str, position: str = "end") -> dict:
    """用 python-docx 编辑 Word 文档（追加文本）"""
    try:
        from docx import Document

        if not Path(filepath).exists():
            return {"success": False, "error": f"文件不存在: {filepath}"}

        doc = Document(filepath)

        if position == "end":
            for line in text.split("\n"):
                if line.strip():
                    doc.add_paragraph(line.strip())
        elif position == "start":
            # 在第一个段落前插入
            for line in reversed(text.split("\n")):
                if line.strip():
                    p = doc.paragraphs[0].insert_paragraph_before(line.strip())

        doc.save(filepath)
        return {"success": True, "engine": "PURE"}
    except Exception as e:
        return {"success": False, "error": str(e)}


def pure_info_word(filepath: str) -> dict:
    """获取 Word 文档基本信息"""
    try:
        from docx import Document

        if not Path(filepath).exists():
            return {"success": False, "error": f"文件不存在: {filepath}"}

        doc = Document(filepath)
        stat = Path(filepath).stat()

        return {
            "success": True,
            "name": Path(filepath).name,
            "path": filepath,
            "size": stat.st_size,
            "paragraph_count": len(doc.paragraphs),
            "engine": "PURE",
        }
    except Exception as e:
        return {"success": False, "error": str(e)}


# ==================== Excel 纯 Python 模式 ====================

def pure_create_excel(name: str, sheets: List[str], filepath: str, data: Dict[str, List[List[Any]]] = None) -> dict:
    """用 openpyxl 创建 Excel 文档"""
    try:
        import openpyxl

        wb = openpyxl.Workbook()
        # 删除默认 sheet
        wb.remove(wb.active)

        for sheet_name in sheets:
            ws = wb.create_sheet(title=sheet_name)
            # 如果有数据
            if data and sheet_name in data:
                for r, row in enumerate(data[sheet_name], start=1):
                    for c, val in enumerate(row, start=1):
                        ws.cell(row=r, column=c, value=val)

        Path(filepath).parent.mkdir(parents=True, exist_ok=True)
        wb.save(filepath)
        return {"success": True, "path": filepath, "sheets": sheets, "engine": "PURE"}
    except Exception as e:
        return {"success": False, "error": str(e)}


def pure_input_excel(filepath: str, sheet: str, data: List[List[Any]]) -> dict:
    """用 openpyxl 录入数据"""
    try:
        import openpyxl

        if not Path(filepath).exists():
            return {"success": False, "error": f"文件不存在: {filepath}"}

        wb = openpyxl.load_workbook(filepath)
        if sheet not in wb.sheetnames:
            ws = wb.create_sheet(title=sheet)
        else:
            ws = wb[sheet]

        for r, row in enumerate(data, start=1):
            for c, val in enumerate(row, start=1):
                ws.cell(row=r, column=c, value=val)

        wb.save(filepath)
        return {"success": True, "engine": "PURE"}
    except Exception as e:
        return {"success": False, "error": str(e)}


def pure_info_excel(filepath: str) -> dict:
    """获取 Excel 文档基本信息"""
    try:
        import openpyxl

        if not Path(filepath).exists():
            return {"success": False, "error": f"文件不存在: {filepath}"}

        wb = openpyxl.load_workbook(filepath, read_only=True)
        stat = Path(filepath).stat()

        return {
            "success": True,
            "name": Path(filepath).name,
            "path": filepath,
            "sheets": wb.sheetnames,
            "size": stat.st_size,
            "sheet_count": len(wb.sheetnames),
            "engine": "PURE",
        }
    except Exception as e:
        return {"success": False, "error": str(e)}


# ==================== PPT 纯 Python 模式（使用 python-pptx）====================

def pure_create_ppt(title: str, filepath: str, slides_content: List[Dict] = None) -> dict:
    """用 python-pptx 创建 PPT"""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt

        prs = Presentation()

        # 标题页
        slide_layout = prs.slide_layouts[0]  # Title layout
        slide = prs.slides.add_slide(slide_layout)
        if slide.shapes.title:
            slide.shapes.title.text = title

        # 添加更多幻灯片
        if slides_content:
            for content in slides_content:
                slide_layout = prs.slide_layouts[1]  # Title and Content
                slide = prs.slides.add_slide(slide_layout)
                if slide.shapes.title:
                    slide.shapes.title.text = content.get("title", "")
                # 添加内容
                for placeholder in slide.placeholders:
                    if placeholder.placeholder_format.idx == 1:  # body
                        placeholder.text = content.get("body", "")

        Path(filepath).parent.mkdir(parents=True, exist_ok=True)
        prs.save(filepath)
        return {"success": True, "path": filepath, "engine": "PURE"}
    except ImportError:
        return {
            "success": False,
            "error": "PPT 纯 Python 模式需要安装 python-pptx：pip install python-pptx",
        }
    except Exception as e:
        return {"success": False, "error": str(e)}


def pure_info_ppt(filepath: str) -> dict:
    """获取 PPT 文档基本信息"""
    try:
        from pptx import Presentation

        if not Path(filepath).exists():
            return {"success": False, "error": f"文件不存在: {filepath}"}

        prs = Presentation(filepath)
        stat = Path(filepath).stat()

        return {
            "success": True,
            "name": Path(filepath).name,
            "path": filepath,
            "size": stat.st_size,
            "slide_count": len(prs.slides),
            "engine": "PURE",
        }
    except ImportError:
        return {"success": False, "error": "需要安装 python-pptx"}
    except Exception as e:
        return {"success": False, "error": str(e)}


# ==================== 安全路径 ====================

def pure_safe_path(path_str: str) -> Path:
    """安全路径校验"""
    path = Path(path_str).resolve()
    parent = path.parent

    if not parent.exists():
        raise FileNotFoundError(f"目录不存在：{parent}")

    # 检查特殊字符
    problematic = '<>"|?*' if __import__("platform").system() == "Windows" else ""
    for ch in problematic:
        if ch in path.name:
            raise ValueError(f"文件名包含非法字符 '{ch}': {path.name}")

    return path


# ==================== 自动选择引擎的 API ====================

def create_word_doc(title: str, filepath: str = "", body: str = "", engine: str = None) -> dict:
    """
    自动选择引擎创建 Word 文档
    - 如果 filepath 已存在且是现有文件 → WPS 模式（编辑）
    - 如果 filepath 不存在 → Pure 模式（新建）
    """
    eng = engine or ENGINE

    if not filepath:
        desktop = Path.home() / "Desktop"
        filepath = str(desktop / f"{title}.docx")

    # 如果文件已存在且引擎是 WPS/MS Office → 用 COM 编辑
    if Path(filepath).exists() and eng in ("WPS", "MSOFFICE"):
        # 调用方负责 COM 操作，这里返回提示
        return {"success": False, "error": "文件已存在，请使用 COM 模式编辑", "hint": "需要 WPS 模式"}

    # 新建文件：纯 Python 即可
    return pure_create_word(title, filepath, body)


def get_engine_info() -> dict:
    """返回当前检测到的引擎信息"""
    return {
        "engine": ENGINE,
        "wps_available": ENGINE == "WPS",
        "msoffice_available": ENGINE == "MSOFFICE",
        "pure_available": True,
    }
