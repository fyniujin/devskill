"""Local Document Generator — generates DOCX/PPX files using python-docx/python-pptx"""
from pathlib import Path
from typing import List, Optional, Dict, Any

try:
    from docx import Document
    from docx.shared import Pt, Inches
    from docx.enum.text import WD_ALIGN_PARAGRAPH
except ImportError:
    Document = None

try:
    from pptx import Presentation
    from pptx.util import Inches as PptxInches, Pt as PptxPt
except ImportError:
    Presentation = None


class DocxGenerator:
    """Generate Word documents using python-docx"""
    
    def __init__(self, template_path: Optional[str] = None):
        if Document is None:
            raise ImportError("python-docx is required. Install with: pip install python-docx")
        self.doc = Document(template_path) if template_path else Document()
    
    def add_heading(self, text: str, level: int = 1):
        self.doc.add_heading(text, level=level)
    
    def add_paragraph(self, text: str, bold: bool = False, size: int = 11):
        p = self.doc.add_paragraph()
        run = p.add_run(text)
        run.font.size = Pt(size)
        run.font.bold = bold
        return p
    
    def add_bullet_list(self, items: List[str], title: Optional[str] = None):
        if title:
            self.add_paragraph(title, bold=True)
        for item in items:
            self.add_paragraph(f"• {item}")
    
    def add_table(self, headers: List[str], rows: List[List[str]]):
        table = self.doc.add_table(rows=1 + len(rows), cols=len(headers))
        for i, h in enumerate(headers):
            table.rows[0].cells[i].text = h
        for r, row in enumerate(rows):
            for c, cell in enumerate(row):
                table.rows[r + 1].cells[c].text = str(cell)
    
    def add_image(self, image_path: str, width: float = 6.0):
        self.doc.add_picture(image_path, width=Inches(width))
    
    def save(self, output_path: str):
        Path(output_path).parent.mkdir(parents=True, exist_ok=True)
        self.doc.save(output_path)
        return output_path


class PptxGenerator:
    """Generate PowerPoint presentations using python-pptx"""
    
    def __init__(self, template_path: Optional[str] = None):
        if Presentation is None:
            raise ImportError("python-pptx is required. Install with: pip install python-pptx")
        self.prs = Presentation(template_path) if template_path else Presentation()
    
    def add_title_slide(self, title: str, subtitle: str = ""):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[0])
        slide.shapes.title.text = title
        if subtitle and len(slide.placeholders) > 1:
            slide.placeholders[1].text = subtitle
    
    def add_content_slide(self, title: str, bullets: List[str]):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        slide.shapes.title.text = title
        text_frame = slide.placeholders[1].text_frame
        for i, bullet in enumerate(bullets):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            p.text = bullet
            p.level = 0
    
    def add_image_slide(self, title: str, image_path: str):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
        slide.shapes.title.text = title
        slide.shapes.add_picture(image_path, PptxInches(1), PptxInches(1.5),
                                width=PptxInches(8), height=PptxInches(5))
    
    def save(self, output_path: str):
        Path(output_path).parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(output_path)
        return output_path


class MindmapGenerator:
    """Generate Mindmap files (export as SVG/PNG)"""
    
    def __init__(self):
        self.nodes = []
        self.edges = []
    
    def add_node(self, node_id: str, label: str, parent_id: Optional[str] = None):
        self.nodes.append({"id": node_id, "label": label})
        if parent_id:
            self.edges.append({"source": parent_id, "target": node_id})
    
    def to_dict(self) -> Dict[str, Any]:
        return {"nodes": self.nodes, "edges": self.edges}
    
    def render_svg(self, output_path: str) -> str:
        """Render a simple SVG mindmap"""
        import math
        Path(output_path).parent.mkdir(parents=True, exist_ok=True)
        width, height = 800, 600
        svg_parts = [
            f'<svg xmlns="http://www.w3.org/2000/svg" width="{width}" height="{height}">',
            '<rect width="100%" height="100%" fill="#ffffff"/>'
        ]
        cx_root, cy_root = width // 2, height // 2
        node_positions = {}
        
        for i, node in enumerate(self.nodes):
            angle = 2 * math.pi * i / max(len(self.nodes), 1)
            x = cx_root + int(200 * math.cos(angle))
            y = cy_root + int(150 * math.sin(angle))
            node_positions[node["id"]] = (x, y)
            svg_parts.append(f'<circle cx="{x}" cy="{y}" r="40" fill="#4A90D9" opacity="0.8"/>')
            svg_parts.append(f'<text x="{x}" y="{y+5}" text-anchor="middle" font-size="12" fill="#ffffff">{node["label"]}</text>')
        
        for edge in self.edges:
            src = node_positions.get(edge["source"])
            tgt = node_positions.get(edge["target"])
            if src and tgt:
                svg_parts.append(f'<line x1="{src[0]}" y1="{src[1]}" x2="{tgt[0]}" y2="{tgt[1]}" stroke="#666" stroke-width="1"/>')
        
        svg_parts.append('</svg>')
        with open(output_path, "w", encoding="utf-8") as f:
            f.write("\n".join(svg_parts))
        return output_path


class FlowchartGenerator:
    """Generate Flowchart SVG from Mermaid-like syntax"""
    
    def __init__(self):
        self.steps = []
    
    def add_step(self, step_id: str, label: str, step_type: str = "process"):
        """
        step_type: start, end, process, decision
        """
        self.steps.append({"id": step_id, "label": label, "type": step_type})
    
    def render_svg(self, output_path: str) -> str:
        Path(output_path).parent.mkdir(parents=True, exist_ok=True)
        width, height = 400, 100 + len(self.steps) * 80
        svg_parts = [
            f'<svg xmlns="http://www.w3.org/2000/svg" width="{width}" height="{height}">',
            '<rect width="100%" height="100%" fill="#ffffff"/>'
        ]
        cx = width // 2
        for i, step in enumerate(self.steps):
            y = 50 + i * 80
            if step["type"] == "start" or step["type"] == "end":
                svg_parts.append(f'<ellipse cx="{cx}" cy="{y}" rx="60" ry="25" fill="#4A90D9" stroke="#2c5f8a" stroke-width="2"/>')
            elif step["type"] == "decision":
                svg_parts.append(f'<polygon points="{cx},{y-25} {cx+60},{y} {cx},{y+25} {cx-60},{y}" fill="#F5A623" stroke="#c17d10" stroke-width="2"/>')
            else:  # process
                svg_parts.append(f'<rect x="{cx-60}" y="{y-25}" width="120" height="50" rx="5" fill="#7ED321" stroke="#4a9e12" stroke-width="2"/>')
            svg_parts.append(f'<text x="{cx}" y="{y+5}" text-anchor="middle" font-size="12" fill="#333333">{step["label"]}</text>')
            if i < len(self.steps) - 1:
                svg_parts.append(f'<line x1="{cx}" y1="{y+30}" x2="{cx}" y2="{y+50}" stroke="#666" stroke-width="2" marker-end="url(#arrow)"/>')
        
        # Arrow marker definition
        svg_parts.insert(2, '<defs><marker id="arrow" markerWidth="10" markerHeight="10" refX="9" refY="3" orient="auto"><path d="M0,0 L0,6 L9,3 z" fill="#666"/></marker></defs>')
        svg_parts.append('</svg>')
        
        with open(output_path, "w", encoding="utf-8") as f:
            f.write("\n".join(svg_parts))
        return output_path
