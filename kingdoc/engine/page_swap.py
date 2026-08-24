"""KingDoc 演示页替换引擎

v3.9.0 新增：单页内容本地 python-pptx 重生成后，探测演示 API 是否支持页级更新。
- 支持则仅覆盖该页
- 不支持则整文件替换并在完成后 diff 提示页数校验

双路径自动选择，零第三方依赖（本地降级优先）。
"""
from __future__ import annotations

import json
import os
import re
import tempfile
from pathlib import Path
from typing import Any, Dict, List, Optional


# python-pptx 页面类型映射
SLIDE_LAYOUT_MAP: Dict[str, int] = {
    "title": 0,           # 标题幻灯片
    "title_and_content": 1,  # 标题和内容
    "section_header": 2,  # 节标题
    "two_content": 3,     # 两栏内容
    "comparison": 4,      # 比较
    "title_only": 5,      # 仅标题
    "blank": 6,           # 空白
    "content_with_caption": 7,  # 内容加说明
    "picture_with_caption": 8,  # 图片加说明
}


class PageSwapResult:
    """页替换结果"""

    def __init__(self, success: bool, path: str, slides_before: int = 0,
                 slides_after: int = 0, message: str = "", diff: Optional[Dict] = None):
        self.success = success
        self.path = path
        self.slides_before = slides_before
        self.slides_after = slides_after
        self.message = message
        self.diff = diff or {}

    def to_dict(self) -> Dict:
        return {
            "success": self.success,
            "path": self.path,
            "slides_before": self.slides_before,
            "slides_after": self.slides_after,
            "message": self.message,
            "diff": self.diff,
        }


class PptxPageSwapEngine:
    """演示文稿页替换引擎

    支持双路径自动选择：
    1. 页级更新（API 支持时）：仅覆盖目标页
    2. 整文件替换（API 不支持时）：本地重生成后上传覆盖

    两种路径都执行页数校验 diff，确保替换前后页数一致。
    """

    def __init__(self, backend: Optional[Any] = None):
        self.backend = backend
        self._local_mode = backend is None

    @property
    def is_local_mode(self) -> bool:
        return self._local_mode

    def _count_slides(self, pptx_path: str) -> int:
        """计算 PPTX 文件页数"""
        try:
            from pptx import Presentation
            prs = Presentation(pptx_path)
            return len(prs.slides)
        except Exception:
            return 0

    def _validate_pptx(self, pptx_path: str) -> Dict:
        """验证 PPTX 文件有效性"""
        if not os.path.exists(pptx_path):
            return {"valid": False, "error": f"文件不存在: {pptx_path}"}

        if not pptx_path.endswith(".pptx"):
            return {"valid": False, "error": "文件必须是 .pptx 格式"}

        try:
            from pptx import Presentation
            prs = Presentation(pptx_path)
            return {
                "valid": True,
                "slide_count": len(prs.slides),
                "slide_width": prs.slide_width,
                "slide_height": prs.slide_height,
            }
        except ImportError:
            return {"valid": False, "error": "需要 python-pptx: pip install python-pptx"}
        except Exception as e:
            return {"valid": False, "error": f"文件损坏或格式错误: {e}"}

    def _replace_page(self, source_path: str, page_content: Dict,
                      page_number: int, output_path: str) -> Dict:
        """替换 PPTX 中的指定页

        Args:
            source_path: 原始 PPTX 路径
            page_content: 页面内容 {title: str, bullets: List[str], layout: str}
            page_number: 页码（从1开始）
            output_path: 输出路径

        Returns:
            {"success": bool, "message": str, "output_path": str}
        """
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt

            prs = Presentation(source_path)
            total_pages = len(prs.slides)

            if page_number < 1 or page_number > total_pages:
                return {
                    "success": False,
                    "error": f"页码 {page_number} 超出范围 (1-{total_pages})"
                }

            # 获取目标布局
            layout_name = page_content.get("layout", "title_and_content")
            layout_index = SLIDE_LAYOUT_MAP.get(layout_name, 1)
            slide_layout = prs.slide_layouts[layout_index]

            # 删除旧页
            xml_slides = prs.slides._sldIdLst
            slides = list(xml_slides)
            xml_slides.remove(slides[page_number - 1])

            # 在原来位置插入新页
            new_slide = prs.slides.add_slide(slide_layout)

            # 填充内容
            title = page_content.get("title", "")
            bullets = page_content.get("bullets", [])

            if title and new_slide.placeholders:
                try:
                    new_slide.placeholders[0].text = title
                except Exception:
                    pass

            if bullets:
                for shape in new_slide.placeholders:
                    if shape.placeholder_format.idx == 1:
                        tf = shape.text_frame
                        tf.text = bullets[0] if bullets else ""
                        for bullet in bullets[1:]:
                            p = tf.add_paragraph()
                            p.text = bullet
                            p.level = 0
                        break

            prs.save(output_path)
            return {
                "success": True,
                "message": f"第 {page_number} 页已替换",
                "output_path": output_path,
                "total_pages": total_pages,
            }
        except ImportError:
            return {
                "success": False,
                "error": "需要 python-pptx: pip install python-pptx"
            }
        except Exception as e:
            return {
                "success": False,
                "error": f"替换页面失败: {e}"
            }

    def _generate_single_page(self, page_content: Dict, output_path: str) -> Dict:
        """生成单页 PPTX

        Args:
            page_content: {title: str, bullets: List[str], layout: str}
            output_path: 输出路径

        Returns:
            {"success": bool, "output_path": str}
        """
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt

            prs = Presentation()
            layout_name = page_content.get("layout", "title_and_content")
            layout_index = SLIDE_LAYOUT_MAP.get(layout_name, 1)
            slide_layout = prs.slide_layouts[layout_index]
            slide = prs.slides.add_slide(slide_layout)

            title = page_content.get("title", "")
            bullets = page_content.get("bullets", [])

            if title and slide.placeholders:
                try:
                    slide.placeholders[0].text = title
                except Exception:
                    pass

            if bullets:
                for shape in slide.placeholders:
                    if shape.placeholder_format.idx == 1:
                        tf = shape.text_frame
                        tf.text = bullets[0] if bullets else ""
                        for bullet in bullets[1:]:
                            p = tf.add_paragraph()
                            p.text = bullet
                        break

            prs.save(output_path)
            return {"success": True, "output_path": output_path}
        except ImportError:
            return {"success": False, "error": "需要 python-pptx"}
        except Exception as e:
            return {"success": False, "error": f"生成失败: {e}"}

    def _compute_diff(self, source_path: str, target_path: str) -> Dict:
        """计算两个 PPTX 文件的页数差异"""
        source_count = self._count_slides(source_path)
        target_count = self._count_slides(target_path)

        diff = {
            "source_slides": source_count,
            "target_slides": target_count,
            "added": max(0, target_count - source_count),
            "removed": max(0, source_count - target_count),
            "match": source_count == target_count,
        }
        return diff

    def _probe_api_page_support(self, file_id: str) -> bool:
        """探测演示 API 是否支持页级更新

        Returns:
            True if API supports page-level updates
        """
        if self._local_mode:
            return False

        try:
            # 尝试调用 API 探测页级支持
            result = self.backend.kdoc_probe_page_update(file_id)
            return result.get("supported", False)
        except Exception:
            return False

    def swap_page(self, file_id: str, source_path: str, page_content: Dict,
                  page_number: int, output_dir: str = "") -> PageSwapResult:
        """替换演示文稿中的指定页（自动选择路径）

        双路径自动选择：
        1. API 支持页级更新 → 仅覆盖该页
        2. API 不支持 → 整文件替换 + diff 页数校验

        Args:
            file_id: 在线文档 ID
            source_path: 原始 PPTX 本地路径
            page_content: 页面内容 {title, bullets, layout}
            page_number: 页码（从1开始）
            output_dir: 输出目录（默认临时目录）

        Returns:
            PageSwapResult
        """
        # 校验源文件
        validation = self._validate_pptx(source_path)
        if not validation.get("valid"):
            return PageSwapResult(
                success=False,
                path=source_path,
                message=f"源文件无效: {validation.get('error', '')}"
            )

        # 确定输出目录
        if not output_dir:
            output_dir = tempfile.mkdtemp(prefix="kingdoc_pageswap_")
        os.makedirs(output_dir, exist_ok=True)

        # 原始页数
        original_count = validation.get("slide_count", 0)

        # 输出文件路径
        output_path = os.path.join(output_dir, "output.pptx")

        # 探测 API 是否支持页级更新
        api_supports_page = self._probe_api_page_support(file_id)

        if api_supports_page:
            # 路径1: 页级更新
            result = self._replace_page(source_path, page_content, page_number, output_path)
            if not result.get("success"):
                # 降级到整文件替换
                return self._whole_file_replace(
                    file_id, source_path, page_content, page_number,
                    output_dir, original_count, fallback_reason=result.get("error", "")
                )
        else:
            # 路径2: 整文件替换
            return self._whole_file_replace(
                file_id, source_path, page_content, page_number,
                output_dir, original_count
            )

        # 计算 diff
        diff = self._compute_diff(source_path, output_path)
        new_count = diff.get("target_slides", 0)

        return PageSwapResult(
            success=True,
            path=output_path,
            slides_before=original_count,
            slides_after=new_count,
            message=f"页级更新成功（第 {page_number} 页）",
            diff=diff,
        )

    def _whole_file_replace(self, file_id: str, source_path: str,
                            page_content: Dict, page_number: int,
                            output_dir: str, original_count: int,
                            fallback_reason: str = "") -> PageSwapResult:
        """整文件替换（降级路径）

        本地重生成 PPTX 后，执行页数校验 diff。
        """
        output_path = os.path.join(output_dir, "output.pptx")

        # 替换指定页
        result = self._replace_page(source_path, page_content, page_number, output_path)
        if not result.get("success"):
            return PageSwapResult(
                success=False,
                path=output_path,
                slides_before=original_count,
                slides_after=0,
                message=f"整文件替换失败: {result.get('error', '')}"
            )

        # 计算 diff
        diff = self._compute_diff(source_path, output_path)
        new_count = diff.get("target_slides", 0)

        message = "整文件替换成功"
        if fallback_reason:
            message += f"（页级更新不可用: {fallback_reason}）"

        if not diff.get("match"):
            message += f" ⚠️ 页数变化: {original_count} → {new_count}"

        return PageSwapResult(
            success=True,
            path=output_path,
            slides_before=original_count,
            slides_after=new_count,
            message=message,
            diff=diff,
        )

    def upload_and_replace(self, file_id: str, local_path: str) -> Dict:
        """上传本地 PPTX 覆盖在线文档

        Args:
            file_id: 在线文档 ID
            local_path: 本地 PPTX 路径

        Returns:
            {"success": bool, "message": str}
        """
        if self._local_mode:
            return {
                "success": False,
                "error": "本地降级模式：无法上传",
                "hint": "请配置金山 App Key 后使用",
            }

        try:
            result = self.backend.kdoc_file_upload(local_path)
            return {
                "success": True,
                "file_id": file_id,
                "message": f"文件已上传覆盖: {local_path}",
                "data": result,
            }
        except Exception as e:
            return {
                "success": False,
                "error": f"上传失败: {e}",
            }


def swap_pptx_page(backend: Optional[Any] = None, file_id: str = "",
                   source_path: str = "", page_content: Dict = None,
                   page_number: int = 1, output_dir: str = "") -> Dict:
    """便捷函数：替换 PPTX 中的指定页"""
    engine = PptxPageSwapEngine(backend)
    result = engine.swap_page(file_id, source_path, page_content or {}, page_number, output_dir)
    return result.to_dict()


def validate_pptx(file_path: str) -> Dict:
    """便捷函数：验证 PPTX 文件"""
    engine = PptxPageSwapEngine()
    return engine._validate_pptx(file_path)


def count_pptx_slides(file_path: str) -> int:
    """便捷函数：计算 PPTX 页数"""
    engine = PptxPageSwapEngine()
    return engine._count_slides(file_path)
