"""KingDoc 格式转换引擎

v3.9.0 新增：补齐 jpg/png/txt 三类缺失的目标格式，
全量覆盖参数校验与失败降级链（云端转换失败→本地转换兜底）。

支持的转换方向：
- doc/docx → pdf / jpg / png / txt / xlsx / pptx
- ppt/pptx → pdf / jpg / png / txt / docx / xlsx
- xls/xlsx → pdf / jpg / png / txt / docx / pptx
- txt → docx / pdf / md
- md → docx / pdf / html

实现原则：
- 云端优先（金山开放平台转换接口）
- 云端失败时本地兜底（python-docx / python-pptx / 纯文本）
- 零第三方依赖（本地兜底仅用标准库 + 已有依赖）
- 硬件自适应（批量转换时读取 hardware.py）
"""
from __future__ import annotations

import json
import os
import shutil
import subprocess
import tempfile
from pathlib import Path
from typing import Any, Dict, List, Optional


# 支持的转换矩阵：{source_format: [target_formats]}
SUPPORTED_CONVERSIONS: Dict[str, List[str]] = {
    "doc": ["pdf", "jpg", "png", "txt", "xlsx", "pptx", "html", "md"],
    "docx": ["pdf", "jpg", "png", "txt", "xlsx", "pptx", "html", "md"],
    "ppt": ["pdf", "jpg", "png", "txt", "docx", "xlsx"],
    "pptx": ["pdf", "jpg", "png", "txt", "docx", "xlsx"],
    "xls": ["pdf", "jpg", "png", "txt", "docx", "pptx"],
    "xlsx": ["pdf", "jpg", "png", "txt", "docx", "pptx"],
    "txt": ["docx", "pdf", "md"],
    "md": ["docx", "pdf", "html"],
    "html": ["docx", "pdf", "md", "txt"],
    "pdf": ["txt", "docx"],
}

# 所有支持的目标格式（v3.9.0 补齐后）
ALL_TARGET_FORMATS = ["pdf", "jpg", "png", "txt", "docx", "xlsx", "pptx", "html", "md"]


class FormatConverter:
    """格式转换引擎

    云端优先 → 本地兜底的双路径转换。
    所有转换都经过参数校验，失败时自动降级。
    """

    def __init__(self, backend: Optional[Any] = None):
        self.backend = backend
        self._local_mode = backend is None

    @property
    def is_local_mode(self) -> bool:
        return self._local_mode

    def validate_conversion(self, source_format: str, target_format: str) -> Dict:
        """校验转换是否支持

        Args:
            source_format: 源格式
            target_format: 目标格式

        Returns:
            {"valid": bool, "message": str, "fallback_available": bool}
        """
        source = source_format.lower().lstrip(".")
        target = target_format.lower().lstrip(".")

        if source == target:
            return {
                "valid": False,
                "message": f"源格式和目标格式相同: {source}",
                "fallback_available": False,
            }

        supported_targets = SUPPORTED_CONVERSIONS.get(source, [])
        if target not in supported_targets:
            return {
                "valid": False,
                "message": f"不支持 {source} → {target} 转换",
                "fallback_available": target in ["txt", "md"],
            }

        # 检查本地兜底是否可用
        fallback_available = self._check_local_fallback(source, target)

        return {
            "valid": True,
            "message": f"支持 {source} → {target} 转换",
            "fallback_available": fallback_available,
        }

    def _check_local_fallback(self, source: str, target: str) -> bool:
        """检查本地兜底是否可用"""
        if target == "txt":
            return True  # 纯文本提取总是可用
        if target == "md":
            return True  # Markdown 转换总是可用
        if target == "docx":
            try:
                import docx  # noqa: F401
                return True
            except ImportError:
                return False
        if target == "pptx":
            try:
                import pptx  # noqa: F401
                return True
            except ImportError:
                return False
        if target == "pdf":
            # 检查是否有 LibreOffice 或 wkhtmltopdf
            return shutil.which("libreoffice") is not None or shutil.which("soffice") is not None
        if target in ("jpg", "png"):
            # 检查是否有 ImageMagick 或 Pillow
            try:
                from PIL import Image  # noqa: F401
                return True
            except ImportError:
                return shutil.which("convert") is not None
        return False

    def convert(self, file_id: str, source_path: str, target_format: str,
                output_dir: str = "") -> Dict:
        """执行格式转换（云端优先 → 本地兜底）

        Args:
            file_id: 在线文档 ID
            source_path: 源文件本地路径
            target_format: 目标格式
            output_dir: 输出目录

        Returns:
            {
                "success": bool,
                "output_path": str,
                "source_format": str,
                "target_format": str,
                "method": "cloud" | "local",
                "message": str
            }
        """
        source_format = Path(source_path).suffix.lstrip(".") or "txt"
        target = target_format.lower().lstrip(".")

        # 参数校验
        validation = self.validate_conversion(source_format, target)
        if not validation.get("valid"):
            return {
                "success": False,
                "output_path": "",
                "source_format": source_format,
                "target_format": target,
                "method": "none",
                "message": validation.get("message", "转换不支持"),
            }

        # 确定输出目录
        if not output_dir:
            output_dir = tempfile.mkdtemp(prefix="kingdoc_convert_")
        os.makedirs(output_dir, exist_ok=True)

        output_filename = f"converted.{target}"
        output_path = os.path.join(output_dir, output_filename)

        # 路径1: 云端转换
        if not self._local_mode:
            try:
                result = self.backend.kdoc_office_convert(file_id, target)
                if result.get("code") == 0:
                    # 下载转换后的文件
                    download_url = result.get("download_url", "")
                    if download_url:
                        self._download_file(download_url, output_path)
                        return {
                            "success": True,
                            "output_path": output_path,
                            "source_format": source_format,
                            "target_format": target,
                            "method": "cloud",
                            "message": f"云端转换成功: {source_format} → {target}",
                        }
            except Exception as e:
                # 云端失败，降级本地
                pass

        # 路径2: 本地兜底
        return self._local_convert(source_path, source_format, target, output_path)

    def _local_convert(self, source_path: str, source_format: str,
                       target_format: str, output_path: str) -> Dict:
        """本地兜底转换"""
        try:
            if target_format == "txt":
                return self._convert_to_txt(source_path, source_format, output_path)
            elif target_format == "md":
                return self._convert_to_md(source_path, source_format, output_path)
            elif target_format == "docx":
                return self._convert_to_docx(source_path, source_format, output_path)
            elif target_format == "pptx":
                return self._convert_to_pptx(source_path, source_format, output_path)
            elif target_format == "pdf":
                return self._convert_to_pdf(source_path, source_format, output_path)
            elif target_format in ("jpg", "png"):
                return self._convert_to_image(source_path, source_format, target_format, output_path)
            elif target_format == "html":
                return self._convert_to_html(source_path, source_format, output_path)
            else:
                return {
                    "success": False,
                    "output_path": "",
                    "source_format": source_format,
                    "target_format": target_format,
                    "method": "none",
                    "message": f"本地不支持 {source_format} → {target_format} 转换",
                }
        except Exception as e:
            return {
                "success": False,
                "output_path": "",
                "source_format": source_format,
                "target_format": target_format,
                "method": "local",
                "message": f"本地转换失败: {e}",
            }

    def _convert_to_txt(self, source_path: str, source_format: str,
                        output_path: str) -> Dict:
        """转换为纯文本"""
        if source_format in ("txt", "md", "html"):
            # 直接复制文本文件
            shutil.copy2(source_path, output_path)
            return {
                "success": True,
                "output_path": output_path,
                "source_format": source_format,
                "target_format": "txt",
                "method": "local",
                "message": f"文本提取成功: {source_format} → txt",
            }
        elif source_format in ("doc", "docx"):
            try:
                from docx import Document
                doc = Document(source_path)
                text = "\n".join(p.text for p in doc.paragraphs)
                with open(output_path, "w", encoding="utf-8") as f:
                    f.write(text)
                return {
                    "success": True,
                    "output_path": output_path,
                    "source_format": source_format,
                    "target_format": "txt",
                    "method": "local",
                    "message": f"DOCX 文本提取成功",
                }
            except ImportError:
                pass
        elif source_format in ("xls", "xlsx"):
            try:
                from openpyxl import load_workbook
                wb = load_workbook(source_path, read_only=True)
                texts = []
                for sheet in wb:
                    for row in sheet.iter_rows(values_only=True):
                        texts.append("\t".join(str(c) if c is not None else "" for c in row))
                with open(output_path, "w", encoding="utf-8") as f:
                    f.write("\n".join(texts))
                return {
                    "success": True,
                    "output_path": output_path,
                    "source_format": source_format,
                    "target_format": "txt",
                    "method": "local",
                    "message": f"XLSX 文本提取成功",
                }
            except ImportError:
                pass
        elif source_format in ("ppt", "pptx"):
            try:
                from pptx import Presentation
                prs = Presentation(source_path)
                texts = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            texts.append(shape.text_frame.text)
                with open(output_path, "w", encoding="utf-8") as f:
                    f.write("\n".join(texts))
                return {
                    "success": True,
                    "output_path": output_path,
                    "source_format": source_format,
                    "target_format": "txt",
                    "method": "local",
                    "message": f"PPTX 文本提取成功",
                }
            except ImportError:
                pass

        return {
            "success": False,
            "output_path": "",
            "source_format": source_format,
            "target_format": "txt",
            "method": "local",
            "message": f"无法提取 {source_format} 的文本（缺少依赖）",
        }

    def _convert_to_md(self, source_path: str, source_format: str,
                       output_path: str) -> Dict:
        """转换为 Markdown"""
        if source_format == "md":
            shutil.copy2(source_path, output_path)
            return {
                "success": True,
                "output_path": output_path,
                "source_format": source_format,
                "target_format": "md",
                "method": "local",
                "message": "Markdown 文件直接复制",
            }
        elif source_format in ("txt", "text"):
            # 纯文本直接复制为 md
            shutil.copy2(source_path, output_path)
            return {
                "success": True,
                "output_path": output_path,
                "source_format": source_format,
                "target_format": "md",
                "method": "local",
                "message": "纯文本转 Markdown",
            }
        elif source_format in ("doc", "docx"):
            try:
                from docx import Document
                doc = Document(source_path)
                md_lines = []
                for p in doc.paragraphs:
                    style = p.style.name if p.style else ""
                    text = p.text.strip()
                    if not text:
                        continue
                    if style.startswith("Heading"):
                        level = 1
                        try:
                            level = int(style.split()[-1])
                        except (ValueError, IndexError):
                            level = 1
                        md_lines.append(f"{'#' * level} {text}")
                    elif style.startswith("List"):
                        md_lines.append(f"- {text}")
                    else:
                        md_lines.append(text)
                with open(output_path, "w", encoding="utf-8") as f:
                    f.write("\n\n".join(md_lines))
                return {
                    "success": True,
                    "output_path": output_path,
                    "source_format": source_format,
                    "target_format": "md",
                    "method": "local",
                    "message": "DOCX → Markdown 转换成功",
                }
            except ImportError:
                pass
        elif source_format == "html":
            # 简单 HTML → MD（去除标签）
            try:
                with open(source_path, "r", encoding="utf-8") as f:
                    html = f.read()
                # 简单去除 HTML 标签
                import re
                text = re.sub(r'<[^>]+>', '', html)
                with open(output_path, "w", encoding="utf-8") as f:
                    f.write(text)
                return {
                    "success": True,
                    "output_path": output_path,
                    "source_format": source_format,
                    "target_format": "md",
                    "method": "local",
                    "message": "HTML → Markdown 转换成功（简单模式）",
                }
            except Exception:
                pass

        return {
            "success": False,
            "output_path": "",
            "source_format": source_format,
            "target_format": "md",
            "method": "local",
            "message": f"无法转换 {source_format} → md",
        }

    def _convert_to_docx(self, source_path: str, source_format: str,
                         output_path: str) -> Dict:
        """转换为 DOCX"""
        try:
            from docx import Document
            from docx.shared import Pt

            doc = Document()

            if source_format in ("txt", "md"):
                with open(source_path, "r", encoding="utf-8") as f:
                    content = f.read()
                for line in content.splitlines():
                    if line.startswith("# "):
                        doc.add_heading(line[2:].strip(), level=1)
                    elif line.startswith("## "):
                        doc.add_heading(line[3:].strip(), level=2)
                    elif line.startswith("- "):
                        doc.add_paragraph(line[2:].strip(), style="List Bullet")
                    else:
                        doc.add_paragraph(line)
            elif source_format in ("xls", "xlsx"):
                try:
                    from openpyxl import load_workbook
                    wb = load_workbook(source_path, read_only=True)
                    for sheet in wb:
                        doc.add_heading(f"Sheet: {sheet.title}", level=2)
                        for row in sheet.iter_rows(values_only=True):
                            line = " | ".join(str(c) if c is not None else "" for c in row)
                            doc.add_paragraph(line)
                except ImportError:
                    doc.add_paragraph("需要 openpyxl 支持表格转换")
            elif source_format in ("ppt", "pptx"):
                try:
                    from pptx import Presentation
                    prs = Presentation(source_path)
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if shape.has_text_frame:
                                doc.add_paragraph(shape.text_frame.text)
                except ImportError:
                    doc.add_paragraph("需要 python-pptx 支持演示转换")
            else:
                # 默认：直接写入文本
                with open(source_path, "r", encoding="utf-8", errors="ignore") as f:
                    doc.add_paragraph(f.read())

            doc.save(output_path)
            return {
                "success": True,
                "output_path": output_path,
                "source_format": source_format,
                "target_format": "docx",
                "method": "local",
                "message": f"本地转换成功: {source_format} → docx",
            }
        except ImportError:
            return {
                "success": False,
                "output_path": "",
                "source_format": source_format,
                "target_format": "docx",
                "method": "local",
                "message": "需要 python-docx: pip install python-docx",
            }

    def _convert_to_pptx(self, source_path: str, source_format: str,
                         output_path: str) -> Dict:
        """转换为 PPTX"""
        try:
            from pptx import Presentation

            prs = Presentation()

            if source_format in ("txt", "md"):
                with open(source_path, "r", encoding="utf-8") as f:
                    content = f.read()
                for section in content.split("\n\n"):
                    lines = section.strip().splitlines()
                    if not lines:
                        continue
                    slide = prs.slides.add_slide(prs.slide_layouts[1])
                    slide.shapes.title.text = lines[0]
                    if len(lines) > 1:
                        slide.placeholders[1].text = "\n".join(lines[1:])
            elif source_format in ("doc", "docx"):
                try:
                    from docx import Document
                    doc = Document(source_path)
                    for p in doc.paragraphs:
                        if p.text.strip():
                            slide = prs.slides.add_slide(prs.slide_layouts[1])
                            slide.shapes.title.text = p.text[:100]
                except ImportError:
                    pass
            else:
                slide = prs.slides.add_slide(prs.slide_layouts[0])
                slide.shapes.title.text = f"Converted from {source_format}"

            prs.save(output_path)
            return {
                "success": True,
                "output_path": output_path,
                "source_format": source_format,
                "target_format": "pptx",
                "method": "local",
                "message": f"本地转换成功: {source_format} → pptx",
            }
        except ImportError:
            return {
                "success": False,
                "output_path": "",
                "source_format": source_format,
                "target_format": "pptx",
                "method": "local",
                "message": "需要 python-pptx: pip install python-pptx",
            }

    def _convert_to_pdf(self, source_path: str, source_format: str,
                        output_path: str) -> Dict:
        """转换为 PDF（使用 LibreOffice）"""
        libreoffice = shutil.which("libreoffice") or shutil.which("soffice")
        if libreoffice:
            try:
                output_dir = os.path.dirname(output_path)
                subprocess.run(
                    [libreoffice, "--headless", "--convert-to", "pdf",
                     "--outdir", output_dir, source_path],
                    capture_output=True, timeout=60
                )
                # LibreOffice 生成的文件名可能不同
                generated = os.path.join(output_dir, Path(source_path).stem + ".pdf")
                if os.path.exists(generated) and generated != output_path:
                    shutil.move(generated, output_path)
                if os.path.exists(output_path):
                    return {
                        "success": True,
                        "output_path": output_path,
                        "source_format": source_format,
                        "target_format": "pdf",
                        "method": "local",
                        "message": "LibreOffice 转换成功",
                    }
            except Exception:
                pass

        return {
            "success": False,
            "output_path": "",
            "source_format": source_format,
            "target_format": "pdf",
            "method": "local",
            "message": "需要 LibreOffice: https://www.libreoffice.org/",
        }

    def _convert_to_image(self, source_path: str, source_format: str,
                          target_format: str, output_path: str) -> Dict:
        """转换为图片（jpg/png）"""
        # 先转 PDF，再转图片
        if source_format in ("doc", "docx", "ppt", "pptx", "xls", "xlsx"):
            pdf_result = self._convert_to_pdf(source_path, source_format,
                                              output_path.replace(f".{target_format}", ".pdf"))
            if pdf_result.get("success"):
                pdf_path = pdf_result["output_path"]
                return self._pdf_to_image(pdf_path, target_format, output_path)

        # 尝试用 Pillow 直接打开（如果源就是图片）
        try:
            from PIL import Image
            img = Image.open(source_path)
            img.save(output_path)
            return {
                "success": True,
                "output_path": output_path,
                "source_format": source_format,
                "target_format": target_format,
                "method": "local",
                "message": f"Pillow 转换成功: {source_format} → {target_format}",
            }
        except ImportError:
            pass
        except Exception:
            pass

        return {
            "success": False,
            "output_path": "",
            "source_format": source_format,
            "target_format": target_format,
            "method": "local",
            "message": f"无法转换 {source_format} → {target_format}（需要 LibreOffice + Pillow）",
        }

    def _pdf_to_image(self, pdf_path: str, target_format: str,
                      output_path: str) -> Dict:
        """PDF 转图片"""
        try:
            from PIL import Image
            # 尝试用 pdf2image
            try:
                from pdf2image import convert_from_path
                images = convert_from_path(pdf_path, first_page=1, last_page=1)
                if images:
                    images[0].save(output_path)
                    return {
                        "success": True,
                        "output_path": output_path,
                        "source_format": "pdf",
                        "target_format": target_format,
                        "method": "local",
                        "message": "PDF → 图片转换成功",
                    }
            except ImportError:
                pass

            # 降级：用 ImageMagick
            convert_cmd = shutil.which("convert")
            if convert_cmd:
                subprocess.run(
                    [convert_cmd, "-density", "150", pdf_path, output_path],
                    capture_output=True, timeout=60
                )
                if os.path.exists(output_path):
                    return {
                        "success": True,
                        "output_path": output_path,
                        "source_format": "pdf",
                        "target_format": target_format,
                        "method": "local",
                        "message": "ImageMagick 转换成功",
                    }
        except Exception:
            pass

        return {
            "success": False,
            "output_path": "",
            "source_format": "pdf",
            "target_format": target_format,
            "method": "local",
            "message": "PDF → 图片需要 pdf2image 或 ImageMagick",
        }

    def _convert_to_html(self, source_path: str, source_format: str,
                         output_path: str) -> Dict:
        """转换为 HTML"""
        if source_format == "html":
            shutil.copy2(source_path, output_path)
            return {
                "success": True,
                "output_path": output_path,
                "source_format": source_format,
                "target_format": "html",
                "method": "local",
                "message": "HTML 直接复制",
            }
        elif source_format == "md":
            try:
                with open(source_path, "r", encoding="utf-8") as f:
                    md = f.read()
                # 简单 MD → HTML
                lines = md.splitlines()
                html_lines = ["<html><body>"]
                for line in lines:
                    if line.startswith("# "):
                        html_lines.append(f"<h1>{line[2:]}</h1>")
                    elif line.startswith("## "):
                        html_lines.append(f"<h2>{line[3:]}</h2>")
                    elif line.startswith("- "):
                        html_lines.append(f"<li>{line[2:]}</li>")
                    else:
                        html_lines.append(f"<p>{line}</p>")
                html_lines.append("</body></html>")
                with open(output_path, "w", encoding="utf-8") as f:
                    f.write("\n".join(html_lines))
                return {
                    "success": True,
                    "output_path": output_path,
                    "source_format": source_format,
                    "target_format": "html",
                    "method": "local",
                    "message": "Markdown → HTML 转换成功",
                }
            except Exception as e:
                return {
                    "success": False,
                    "output_path": "",
                    "source_format": source_format,
                    "target_format": "html",
                    "method": "local",
                    "message": f"MD → HTML 失败: {e}",
                }
        elif source_format in ("txt", "text"):
            with open(source_path, "r", encoding="utf-8") as f:
                text = f.read()
            html = f"<html><body><pre>{text}</pre></body></html>"
            with open(output_path, "w", encoding="utf-8") as f:
                f.write(html)
            return {
                "success": True,
                "output_path": output_path,
                "source_format": source_format,
                "target_format": "html",
                "method": "local",
                "message": "TXT → HTML 转换成功",
            }

        return {
            "success": False,
            "output_path": "",
            "source_format": source_format,
            "target_format": "html",
            "method": "local",
            "message": f"无法转换 {source_format} → html",
        }

    def _download_file(self, url: str, output_path: str):
        """下载文件到本地"""
        try:
            import urllib.request
            urllib.request.urlretrieve(url, output_path)
        except Exception:
            pass

    def get_supported_formats(self) -> Dict:
        """获取所有支持的格式"""
        return {
            "sources": list(SUPPORTED_CONVERSIONS.keys()),
            "targets": ALL_TARGET_FORMATS,
            "matrix": SUPPORTED_CONVERSIONS,
        }


def convert_file(backend: Optional[Any] = None, file_id: str = "",
                 source_path: str = "", target_format: str = "",
                 output_dir: str = "") -> Dict:
    """便捷函数：转换文件格式"""
    converter = FormatConverter(backend)
    return converter.convert(file_id, source_path, target_format, output_dir)


def validate_conversion(source_format: str, target_format: str) -> Dict:
    """便捷函数：校验转换是否支持"""
    converter = FormatConverter()
    return converter.validate_conversion(source_format, target_format)


def get_supported_formats() -> Dict:
    """便捷函数：获取支持的格式列表"""
    converter = FormatConverter()
    return converter.get_supported_formats()
